from io import BytesIO
import html
import json
import os
import re
import shutil
import time
import urllib.error
import urllib.request
from pathlib import Path
from uuid import uuid4
from urllib.parse import quote

import pandas as pd
import streamlit as st
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from customer_pdf_generator import build_pdf_draft_from_proposal, customer_pdf_filename, generate_customer_pdf

from poi_importer import (
    FIELD_LABELS as IMPORT_FIELD_LABELS,
    analyze_poi_file,
    merge_poi_data,
)


# HZH brand palette, sampled from the supplied HangZhou Health logo.
NAVY = RGBColor(14, 67, 110)
DEEP_GREEN = RGBColor(24, 181, 178)
GOLD = RGBColor(29, 145, 219)
BEIGE = RGBColor(247, 252, 253)
LIGHT_BEIGE = RGBColor(236, 249, 250)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(45, 55, 65)
MUTED = RGBColor(100, 111, 120)
PLACEHOLDER_FILL = RGBColor(238, 247, 250)
PLACEHOLDER_LINE = RGBColor(162, 213, 225)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BASE_DIR = Path(__file__).parent
STORAGE_DIR = Path(os.getenv("HZH_STORAGE_DIR", str(BASE_DIR))).expanduser()
DATA_DIR = Path(os.getenv("HZH_DATA_DIR", str(STORAGE_DIR / "data"))).expanduser()
OUTPUTS_DIR = Path(os.getenv("HZH_OUTPUTS_DIR", str(STORAGE_DIR / "outputs"))).expanduser()
ASSETS_DIR = BASE_DIR / "assets"
LOGO_DIR = ASSETS_DIR / "logo"
POI_IMAGE_DIR = Path(os.getenv("HZH_POI_IMAGE_DIR", str(STORAGE_DIR / "assets" / "poi_images"))).expanduser()
LOGO_MARK_PATH = LOGO_DIR / "hzh-logo-mark.png"
LOGO_FULL_PATH = LOGO_DIR / "hzh-logo-full.png"
POI_DB_PATH = DATA_DIR / "poi_database.csv"
POI_TEMPLATE_PATH = OUTPUTS_DIR / "poi_template.csv"

# Language rules: keep the internal tool in Chinese and generated client files in English.
interfaceLanguage = "zh-CN"
outputLanguage = "en"

generationStatusLabels = {
    "idle": "未生成",
    "generating": "生成中",
    "completed": "已完成",
    "failed": "生成失败",
    "stale": "左侧信息已修改，请重新生成",
}

POI_FIELD_LABELS = {
    "id": "内部ID",
    "poi_id": "点位编号",
    "code": "点位编码",
    "name_cn": "中文名称",
    "name_en": "英文名称",
    "category": "点位类型",
    "city": "城市",
    "district": "区域",
    "duration_hours": "建议游览时长（小时）",
    "price_min": "最低成本价",
    "price_max": "最高成本价",
    "currency": "币种",
    "tags": "标签",
    "suitable_for": "适合人群",
    "description": "综合描述",
    "description_cn": "中文介绍",
    "description_en": "英文介绍",
    "recommended_reason": "推荐理由",
    "highlights": "亮点",
    "target_users": "目标客群",
    "image_placeholder": "图片占位符或图片路径",
    "image_notes": "图片备注",
    "image_source": "图片来源",
    "image_alt": "图片英文说明",
    "image_url": "网络图片地址",
    "image_path": "本地图片路径",
    "images": "图片组",
    "address": "地址",
    "opening_hours": "开放时间",
    "reservation_info": "预约信息",
    "notes": "内部备注",
    "internal_remark": "内部备注补充",
    "status": "状态",
}

ITINERARY_COLUMN_LABELS = {
    "Day": "天数",
    "City": "英文城市",
    "Theme": "英文主题",
    "Arrangement": "英文行程安排",
}

HIGHLIGHT_COLUMN_LABELS = {
    "Title": "英文亮点标题",
    "Caption": "英文短说明",
}

POI_COLUMNS = [
    "id",
    "poi_id",
    "code",
    "name_cn",
    "name_en",
    "category",
    "city",
    "district",
    "duration_hours",
    "price_min",
    "price_max",
    "currency",
    "tags",
    "suitable_for",
    "description",
    "description_cn",
    "description_en",
    "recommended_reason",
    "highlights",
    "target_users",
    "image_placeholder",
    "image_notes",
    "image_source",
    "image_alt",
    "image_url",
    "image_path",
    "images",
    "address",
    "opening_hours",
    "reservation_info",
    "notes",
    "internal_remark",
    "status",
]

DEFAULT_POIS = [
    {
        "poi_id": "POI-001",
        "name_cn": "国际口腔医疗中心",
        "name_en": "International Dental Care Center",
        "category": "医疗机构",
        "city": "Hangzhou",
        "district": "Sukhumvit",
        "duration_hours": 2,
        "price_min": 800,
        "price_max": 1500,
        "currency": "CNY",
        "tags": "医疗;室内;高端;适合外宾",
        "suitable_for": "康养客户;家庭客户",
        "description_cn": "适合安排口腔咨询、影像检查与治疗方案沟通。",
        "description_en": "Suitable for dental consultation, imaging review, and treatment planning support.",
        "image_placeholder": "Dental clinic image placeholder",
        "image_path": "",
        "image_url": "",
        "image_alt": "",
        "image_source": "",
        "image_notes": "",
        "address": "Hangzhou",
        "notes": "需提前确认医生档期。",
    },
    {
        "poi_id": "POI-002",
        "name_cn": "滨海康养酒店",
        "name_en": "Seaside Wellness Resort",
        "category": "酒店",
        "city": "Shanghai",
        "district": "Patong",
        "duration_hours": 3,
        "price_min": 1200,
        "price_max": 2200,
        "currency": "CNY",
        "tags": "康养;轻体力;高端;亲子",
        "suitable_for": "康养客户;家庭客户;商务客户",
        "description_cn": "适合术后休整、海边轻松活动与高端住宿体验。",
        "description_en": "A recovery-friendly resort option for premium stays and relaxed seaside time.",
        "image_placeholder": "Wellness resort image placeholder",
        "image_path": "",
        "image_url": "",
        "image_alt": "",
        "image_source": "",
        "image_notes": "",
        "address": "Shanghai",
        "notes": "旺季价格波动较大。",
    },
    {
        "poi_id": "POI-003",
        "name_cn": "城市夜游河畔体验",
        "name_en": "Riverside Night Tour",
        "category": "夜游",
        "city": "Hangzhou",
        "district": "Riverside",
        "duration_hours": 2,
        "price_min": 300,
        "price_max": 600,
        "currency": "CNY",
        "tags": "文化;夜间;轻体力;适合外宾",
        "suitable_for": "商务客户;家庭客户",
        "description_cn": "适合作为晚间轻松活动，展示城市夜景与本地文化氛围。",
        "description_en": "A light evening experience featuring city views and local cultural atmosphere.",
        "image_placeholder": "Night tour image placeholder",
        "image_path": "",
        "image_url": "",
        "image_alt": "",
        "image_source": "",
        "image_notes": "",
        "address": "Hangzhou",
        "notes": "雨季建议准备备选室内活动。",
    },
]


CITY_ALIASES = {
    "hangzhou": ["hangzhou", "杭州", "杭州径山"],
    "yiwu": ["yiwu", "义乌"],
    "shanghai": ["shanghai", "上海"],
    "guangzhou": ["guangzhou", "广州", "廣州"],
    "general": ["general", "通用"],
}


def city_matches(poi_city, itinerary_city):
    """Match English itinerary cities to Chinese POI city values."""
    poi_text = clean_text(poi_city).lower()
    city_text = clean_text(itinerary_city).lower()
    if not city_text:
        return False
    if poi_text == city_text:
        return True
    aliases = CITY_ALIASES.get(city_text, [city_text])
    return any(alias.lower() == poi_text or alias.lower() in poi_text for alias in aliases)


def city_key(value):
    """Return the canonical city key used for hard POI matching."""
    text = clean_text(value).lower()
    if not text:
        return ""
    for key, aliases in CITY_ALIASES.items():
        if key == "general":
            continue
        if text == key or any(alias.lower() == text or alias.lower() in text for alias in aliases):
            return key
    return text


def poi_allowed_for_city(poi, itinerary_city):
    """Hard city constraint for route generation and default manual choices."""
    return bool(clean_text(itinerary_city)) and city_matches(poi.get("city"), itinerary_city)


def normalize_itinerary_and_pois(itinerary_rows, day_poi_ids=None, poi_df=None):
    """Keep Day numbers continuous and remove POIs that do not match the day's city."""
    rows = []
    normalized_ids = []
    source_ids = day_poi_ids or []
    poi_lookup = {}
    if poi_df is not None and hasattr(poi_df, "iterrows"):
        poi_lookup = {str(row.get("poi_id")): row.to_dict() for _, row in poi_df.iterrows()}

    for index, row in enumerate(itinerary_records(itinerary_rows), start=1):
        updated = dict(row)
        updated["Day"] = index
        rows.append(updated)
        saved_ids = source_ids[index - 1] if index - 1 < len(source_ids) else []
        city = clean_text(updated.get("City"))
        kept = []
        for poi_id in saved_ids:
            poi = poi_lookup.get(str(poi_id))
            if poi is None or poi_allowed_for_city(poi, city):
                kept.append(poi_id)
            else:
                print(f"POI city mismatch: itinerary city {city} but POI city {poi.get('city')} ({poi_id})")
        normalized_ids.append(kept)
    return rows, normalized_ids


def is_medical_theme(theme):
    text = clean_text(theme).lower()
    return any(word in text for word in ["medical", "health check", "consultation", "screening", "体检", "医疗"])


def validate_generated_route(itinerary_rows, day_poi_ids, poi_df, requirements=None):
    """Post-generation guardrails for day numbering, city matching, and duplicate medical days."""
    rows, normalized_ids = normalize_itinerary_and_pois(itinerary_rows, day_poi_ids, poi_df)
    medical_indices = [index for index, row in enumerate(rows) if is_medical_theme(row.get("Theme"))]
    if len(medical_indices) > 1:
        primary = medical_indices[0]
        for duplicate in medical_indices[1:]:
            city = clean_text(rows[duplicate].get("City"))
            rows[duplicate]["Theme"] = generated_theme(rows[duplicate]["Day"], len(rows), city)
            if is_medical_theme(rows[duplicate]["Theme"]):
                rows[duplicate]["Theme"] = "Wellness & Local Experience"
            rows[duplicate]["Arrangement"] = generic_client_arrangement(rows[duplicate]["Day"], len(rows), city, rows[duplicate]["Theme"])
            print(f"Duplicate medical day adjusted: Day {duplicate + 1} merged into Day {primary + 1}")
        city = clean_text(rows[primary].get("City"))
        rows[primary]["Theme"] = "Group Health Checkup & Recovery"
        rows[primary]["Arrangement"] = f"Coordinated health checkups for all travelers in {city}, with bilingual support, private transfers, and recovery-friendly arrangements."
    return rows, normalized_ids


def reorder_itinerary_by_labels(itinerary_rows, day_poi_ids, sorted_labels):
    """Apply drag-sort labels back to itinerary rows and selected POIs."""
    rows = itinerary_records(itinerary_rows)
    if not rows or not sorted_labels:
        return rows, day_poi_ids
    label_to_index = {itinerary_sort_label(row, index): index for index, row in enumerate(rows)}
    ordered_indices = [label_to_index[label] for label in sorted_labels if label in label_to_index]
    if len(ordered_indices) != len(rows):
        return rows, day_poi_ids
    reordered_rows = [dict(rows[index]) for index in ordered_indices]
    reordered_pois = [day_poi_ids[index] if index < len(day_poi_ids) else [] for index in ordered_indices]
    return normalize_itinerary_and_pois(reordered_rows, reordered_pois, active_poi_df if "active_poi_df" in globals() else None)


def itinerary_sort_label(row, index):
    city = clean_text(row.get("City")) or "City TBD"
    theme = clean_text(row.get("Theme")) or "Untitled"
    return f"Day {index + 1} | {city} | {theme}"



def inject_app_theme():
    """Apply a calm blue visual system to the Streamlit interface."""
    st.markdown(
        """
        <style>
        :root {
          --color-primary: #1E40AF;
          --color-primary-hover: #1D4ED8;
          --color-primary-text: #FFFFFF;
          --color-tag-bg: #EAF2FF;
          --color-tag-text: #1E40AF;
          --color-tag-border: #BFDBFE;
          --color-tag-close: #1E40AF;
          --color-secondary-bg: #FFFFFF;
          --color-secondary-text: #334155;
          --color-secondary-border: #CBD5E1;
          --color-secondary-hover-bg: #F1F5F9;
          --color-danger: #DC2626;
          --color-danger-hover: #B91C1C;
          --color-danger-bg: #FEF2F2;
          --color-muted-bg: #F8FAFC;
          --color-border: #E2E8F0;
          --color-status-idle-bg: #F8FAFC;
          --color-status-idle-text: #64748B;
          --color-status-generating-bg: #EFF6FF;
          --color-status-generating-text: #1D4ED8;
          --color-status-completed-bg: #ECFDF5;
          --color-status-completed-text: #047857;
          --color-status-stale-bg: #FFF7ED;
          --color-status-stale-text: #C2410C;
          --color-status-failed-bg: #FEF2F2;
          --color-status-failed-text: #DC2626;
        }

        /* Primary buttons: Streamlit versions differ, so cover current and older DOM shapes. */
        .stButton button[kind="primary"],
        .stDownloadButton button[kind="primary"],
        button[kind="primary"],
        button[data-testid="stBaseButton-primary"],
        [data-testid="stBaseButton-primary"],
        [data-testid="stButton"] button[data-testid="stBaseButton-primary"],
        [data-testid="stDownloadButton"] button[data-testid="stBaseButton-primary"] {
          background-color: var(--color-primary) !important;
          background: var(--color-primary) !important;
          color: var(--color-primary-text) !important;
          border-color: var(--color-primary) !important;
          box-shadow: 0 1px 2px rgba(15, 23, 42, 0.08) !important;
        }
        .stButton button[kind="primary"]:hover,
        .stDownloadButton button[kind="primary"]:hover,
        button[kind="primary"]:hover,
        button[data-testid="stBaseButton-primary"]:hover,
        [data-testid="stBaseButton-primary"]:hover {
          background-color: var(--color-primary-hover) !important;
          background: var(--color-primary-hover) !important;
          border-color: var(--color-primary-hover) !important;
          color: var(--color-primary-text) !important;
        }
        .stButton button[kind="primary"] p,
        .stDownloadButton button[kind="primary"] p,
        button[data-testid="stBaseButton-primary"] p,
        [data-testid="stBaseButton-primary"] p {
          color: var(--color-primary-text) !important;
        }

        /* Secondary buttons stay low-emphasis. */
        .stButton button:not([kind="primary"]):not([data-testid="stBaseButton-primary"]),
        .stDownloadButton button:not([kind="primary"]):not([data-testid="stBaseButton-primary"]),
        button[data-testid="stBaseButton-secondary"] {
          background-color: var(--color-secondary-bg) !important;
          background: var(--color-secondary-bg) !important;
          color: var(--color-secondary-text) !important;
          border: 1px solid var(--color-secondary-border) !important;
        }
        .stButton button:not([kind="primary"]):not([data-testid="stBaseButton-primary"]):hover,
        .stDownloadButton button:not([kind="primary"]):not([data-testid="stBaseButton-primary"]):hover,
        button[data-testid="stBaseButton-secondary"]:hover {
          background-color: var(--color-secondary-hover-bg) !important;
          background: var(--color-secondary-hover-bg) !important;
          color: var(--color-primary) !important;
          border-color: var(--color-primary) !important;
        }

        /* Multi-select selected tags: light blue, not button blue. */
        [data-baseweb="tag"],
        [data-baseweb="tag"] > span,
        div[data-baseweb="tag"] {
          background-color: var(--color-tag-bg) !important;
          border-color: var(--color-tag-border) !important;
          color: var(--color-tag-text) !important;
        }
        [data-baseweb="tag"] span,
        [data-baseweb="tag"] svg,
        [data-baseweb="tag"] path {
          color: var(--color-tag-text) !important;
          fill: var(--color-tag-close) !important;
        }
        [data-baseweb="select"] [aria-selected="true"],
        [data-baseweb="menu"] li[aria-selected="true"] {
          background-color: var(--color-tag-bg) !important;
          color: var(--color-tag-text) !important;
        }
        [data-baseweb="checkbox"] [aria-checked="true"],
        [data-baseweb="radio"] [aria-checked="true"] {
          border-color: var(--color-primary) !important;
          background-color: var(--color-primary) !important;
        }
        .poi-source-badges {
          display: flex;
          flex-wrap: wrap;
          gap: 8px;
          margin: 6px 0 16px 0;
        }
        .poi-source-badge {
          display: inline-flex;
          align-items: center;
          border-radius: 999px;
          padding: 5px 10px;
          font-size: 13px;
          line-height: 1.2;
          border: 1px solid transparent;
          white-space: nowrap;
        }
        .poi-source-badge.db {
          background: #EAF2FF;
          color: #1E40AF;
          border-color: #BFDBFE;
        }
        .poi-source-badge.ai {
          background: #FFF7ED;
          color: #C2410C;
          border-color: #FDBA74;
        }
        .poi-source-legend {
          color: #64748B;
          font-size: 13px;
          margin: -2px 0 8px 0;
        }
        [data-testid="stSidebar"] {
          border-right: 1px solid var(--color-border);
        }
        [data-testid="stExpander"] {
          border-color: var(--color-border) !important;
          background: #FFFFFF !important;
        }
        [data-testid="stDataFrame"],
        [data-testid="stDataEditor"] {
          border-color: var(--color-border) !important;
        }
        .stButton button {
          min-height: 34px !important;
          border-radius: 8px !important;
          padding: 0.35rem 0.8rem !important;
          white-space: nowrap !important;
        }
        .stButton button p {
          white-space: nowrap !important;
          word-break: keep-all !important;
        }
        button[data-testid="stBaseButton-tertiary"],
        [data-testid="stBaseButton-tertiary"] {
          background-color: var(--color-danger-bg) !important;
          color: var(--color-danger) !important;
          border: 1px solid #FECACA !important;
        }
        button[data-testid="stBaseButton-tertiary"]:hover,
        [data-testid="stBaseButton-tertiary"]:hover {
          background-color: #FEE2E2 !important;
          color: var(--color-danger-hover) !important;
          border-color: var(--color-danger) !important;
        }
        button[data-testid="stBaseButton-tertiary"] p,
        [data-testid="stBaseButton-tertiary"] p {
          color: var(--color-danger) !important;
        }
        .hzh-app-tabs {
          display: flex;
          gap: 32px;
          border-bottom: 1px solid var(--color-border);
          margin: 0.75rem 0 2rem;
        }
        .hzh-app-tab {
          position: relative;
          display: inline-flex;
          align-items: center;
          padding: 0.85rem 0;
          color: #64748B !important;
          text-decoration: none !important;
          font-size: 1.05rem;
          font-weight: 600;
          line-height: 1.2;
        }
        .hzh-app-tab:hover {
          color: var(--color-primary) !important;
        }
        .hzh-app-tab.active {
          color: var(--color-primary) !important;
          font-weight: 700;
        }
        .hzh-app-tab.active::after {
          content: "";
          position: absolute;
          left: 0;
          right: 0;
          bottom: -1px;
          height: 3px;
          background: var(--color-primary);
          border-radius: 999px;
        }
        .hzh-edit-title-row {
          display: flex;
          align-items: center;
          gap: 0.75rem;
          margin: 0.75rem 0 0.15rem;
        }
        .hzh-edit-title-row h1 {
          margin: 0 !important;
        }
        .hzh-poi-list-note {
          color: #64748B;
          font-size: 0.88rem;
          margin-bottom: 0.5rem;
        }
        .hzh-poi-table-wrapper {
          width: 100%;
          overflow-x: auto;
          border: 1px solid var(--color-border);
          border-radius: 10px;
          background: #FFFFFF;
        }
        .hzh-poi-table {
          min-width: 1320px;
          width: 100%;
          border-collapse: collapse;
          table-layout: fixed;
          font-size: 0.9rem;
        }
        .hzh-poi-table th {
          background: #F8FAFC;
          color: #0F172A;
          font-weight: 700;
          text-align: left;
          border-bottom: 1px solid var(--color-border);
          padding: 0.72rem 0.75rem;
          white-space: nowrap;
        }
        .hzh-poi-table td {
          border-bottom: 1px solid var(--color-border);
          color: #334155;
          padding: 0.68rem 0.75rem;
          vertical-align: middle;
          overflow: hidden;
          text-overflow: ellipsis;
        }
        .hzh-poi-table tr:last-child td {
          border-bottom: 0;
        }
        .hzh-poi-table .nowrap {
          white-space: nowrap;
        }
        .hzh-poi-table .muted {
          color: #94A3B8;
        }
        .hzh-poi-action-cell {
          white-space: nowrap;
          min-width: 170px;
        }
        .hzh-poi-action-btn {
          display: inline-flex;
          align-items: center;
          justify-content: center;
          min-width: 64px;
          height: 32px;
          border-radius: 8px;
          padding: 0 0.78rem;
          margin-right: 8px;
          text-decoration: none !important;
          font-weight: 600;
          line-height: 1;
          box-sizing: border-box;
        }
        .hzh-poi-action-btn.edit {
          color: var(--color-primary) !important;
          background: var(--color-tag-bg);
          border: 1px solid var(--color-tag-border);
        }
        .hzh-poi-action-btn.edit:hover {
          background: #DBEAFE;
          border-color: var(--color-primary);
        }
        .hzh-poi-action-btn.delete {
          color: var(--color-danger) !important;
          background: var(--color-danger-bg);
          border: 1px solid #FECACA;
        }
        .hzh-poi-action-btn.delete:hover {
          background: #FEE2E2;
          border-color: var(--color-danger);
        }
        .hzh-poi-col-id { width: 120px; }
        .hzh-poi-col-cn { width: 190px; }
        .hzh-poi-col-en { width: 220px; }
        .hzh-poi-col-type { width: 135px; }
        .hzh-poi-col-city { width: 110px; }
        .hzh-poi-col-duration { width: 85px; }
        .hzh-poi-col-price { width: 185px; }
        .hzh-poi-col-image { width: 125px; }
        .hzh-poi-col-action { width: 180px; }
        .hzh-poi-row-divider {
          border-bottom: 1px solid var(--color-border);
          margin: 0.2rem 0 0.35rem;
        }
        .hzh-status-card {
          border-radius: 10px;
          border: 1px solid var(--color-border);
          padding: 0.85rem 1rem;
          margin: 0.35rem 0 1rem;
          font-size: 0.95rem;
          line-height: 1.45;
        }
        .hzh-status-card strong {
          display: block;
          margin-bottom: 0.25rem;
          font-weight: 700;
        }
        .hzh-status-idle { background: var(--color-status-idle-bg); color: var(--color-status-idle-text); }
        .hzh-status-generating { background: var(--color-status-generating-bg); color: var(--color-status-generating-text); border-color: #BFDBFE; }
        .hzh-status-completed { background: var(--color-status-completed-bg); color: var(--color-status-completed-text); border-color: #A7F3D0; }
        .hzh-status-stale { background: var(--color-status-stale-bg); color: var(--color-status-stale-text); border-color: #FDBA74; }
        .hzh-status-failed { background: var(--color-status-failed-bg); color: var(--color-status-failed-text); border-color: #FECACA; }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_status_card(status, title, message):
    """Render generation status with controlled design-system colors."""
    safe_status = status if status in {"idle", "generating", "completed", "stale", "failed"} else "idle"
    st.markdown(
        f'<div class="hzh-status-card hzh-status-{safe_status}"><strong>{clean_text(title)}</strong>{clean_text(message)}</div>',
        unsafe_allow_html=True,
    )


def money(value, currency="$"):
    """Format a number as a client-facing amount."""
    try:
        amount = float(value)
    except (TypeError, ValueError):
        amount = 0
    if currency == "CNY":
        return f"CNY {amount:,.0f}"
    if currency == "USD" or currency == "$":
        return f"${amount:,.0f}"
    return f"{currency} {amount:,.0f}"


def split_lines(text):
    """Turn a multiline text box into clean bullet items."""
    return [line.strip() for line in str(text or "").splitlines() if line.strip()]


def itinerary_records(itinerary):
    """Accept Streamlit editor output as either a list of rows or a dataframe."""
    if hasattr(itinerary, "to_dict"):
        return itinerary.to_dict("records")
    return itinerary or []


def clean_text(value):
    """Keep PowerPoint text safe and predictable."""
    if pd.isna(value):
        return ""
    return str(value or "").strip()


def optional_number(value):
    """Return a float when present, otherwise None."""
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    if pd.isna(number):
        return None
    return number


def uploaded_file_bytes(uploaded_file):
    """Return reusable bytes from a Streamlit upload object."""
    if not uploaded_file:
        return None
    return uploaded_file.getvalue()


def image_size(image_bytes):
    """Read image dimensions without saving the uploaded file."""
    with Image.open(BytesIO(image_bytes)) as image:
        return image.size


def safe_filename_slug(value):
    """Create an ASCII-safe filename slug for uploaded POI images."""
    import re

    text = clean_text(value).lower()
    text = re.sub(r"[^a-z0-9]+", "_", text)
    return text.strip("_")[:42] or "poi_image"


def save_uploaded_poi_image(uploaded_file, poi_id, name_hint=""):
    """Persist a POI image under assets/poi_images and return a project-relative path."""
    if not uploaded_file:
        return ""
    POI_IMAGE_DIR.mkdir(parents=True, exist_ok=True)
    suffix = Path(uploaded_file.name).suffix.lower()
    if suffix not in [".jpg", ".jpeg", ".png", ".webp"]:
        suffix = ".jpg"
    safe_id = safe_filename_slug(poi_id)
    safe_name = safe_filename_slug(name_hint)
    filename = f"{safe_id}_{safe_name}_{int(time.time())}{suffix}"
    output_path = POI_IMAGE_DIR / filename
    output_path.write_bytes(uploaded_file.getvalue())
    try:
        return str(output_path.relative_to(BASE_DIR))
    except ValueError:
        return str(output_path)


def resolve_local_image_path(path_value):
    """Resolve a stored local image path to an existing filesystem path."""
    raw_path = clean_text(path_value)
    if not raw_path:
        return None
    image_path = Path(raw_path)
    if not image_path.is_absolute():
        image_path = BASE_DIR / raw_path
    if image_path.exists() and image_path.suffix.lower() in [".png", ".jpg", ".jpeg", ".webp"]:
        return image_path
    return None


def poi_has_image(poi):
    """Return whether a POI has a usable local path or a recorded image URL."""
    return bool(resolve_local_image_path(poi.get("image_path")) or clean_text(poi.get("image_url")))


def active_pois_only(df):
    """Return non-deleted POIs for route generation and normal selection."""
    if "status" not in df.columns:
        return df.copy()
    return df[df["status"].astype(str).str.lower() != "deleted"].copy()


def poi_status_label(value):
    return "已删除" if clean_text(value).lower() == "deleted" else "启用"


def build_images_json(image_path="", image_url="", name="", is_main=True):
    url = clean_text(image_path) or clean_text(image_url)
    if not url:
        return "[]"
    image_name = clean_text(name) or Path(url).name
    return json.dumps(
        [
            {
                "id": f"IMG-{uuid4().hex[:8].upper()}",
                "url": url,
                "name": image_name,
                "isMain": bool(is_main),
                "uploadedAt": time.strftime("%Y-%m-%d %H:%M:%S"),
            }
        ],
        ensure_ascii=False,
    )


def update_poi_row(poi_df, poi_id, values):
    updated_df = poi_df.copy()
    row_mask = updated_df["poi_id"].astype(str) == str(poi_id)
    for key, value in values.items():
        if key not in updated_df.columns:
            updated_df[key] = ""
        updated_df.loc[row_mask, key] = value
    return updated_df


def soft_delete_pois(poi_df, poi_ids):
    updated_df = poi_df.copy()
    if "status" not in updated_df.columns:
        updated_df["status"] = "active"
    mask = updated_df["poi_id"].astype(str).isin([str(poi_id) for poi_id in poi_ids])
    updated_df.loc[mask, "status"] = "deleted"
    return updated_df


def restore_pois(poi_df, poi_ids):
    updated_df = poi_df.copy()
    if "status" not in updated_df.columns:
        updated_df["status"] = "active"
    mask = updated_df["poi_id"].astype(str).isin([str(poi_id) for poi_id in poi_ids])
    updated_df.loc[mask, "status"] = "active"
    return updated_df


def normalize_poi_df(df):
    """Ensure imported or edited POI data always has the expected columns."""
    df = df.copy()
    for column in POI_COLUMNS:
        if column not in df.columns:
            df[column] = ""
    df = df[POI_COLUMNS]
    df = df.fillna("")

    for column in ["duration_hours", "price_min", "price_max"]:
        df[column] = pd.to_numeric(df[column], errors="coerce")

    df["currency"] = df["currency"].replace("", "CNY")
    missing_id = df["poi_id"].astype(str).str.strip() == ""
    code_as_id = missing_id & (df["code"].astype(str).str.strip() != "")
    df.loc[code_as_id, "poi_id"] = df.loc[code_as_id, "code"]
    missing_id = df["poi_id"].astype(str).str.strip() == ""
    df.loc[missing_id, "poi_id"] = [f"POI-{uuid4().hex[:8].upper()}" for _ in range(missing_id.sum())]
    df["id"] = df["id"].astype(str).str.strip()
    missing_internal_id = df["id"] == ""
    df.loc[missing_internal_id, "id"] = df.loc[missing_internal_id, "poi_id"]
    df["code"] = df["code"].astype(str).str.strip()
    missing_code = df["code"] == ""
    df.loc[missing_code, "code"] = df.loc[missing_code, "poi_id"]
    df["status"] = df["status"].apply(lambda value: "deleted" if clean_text(value).lower() == "deleted" else "active")
    missing_images = df["images"].astype(str).str.strip() == ""
    df.loc[missing_images, "images"] = "[]"
    return df


def ensure_poi_files():
    """Create the local POI database and import template if they do not exist."""
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
    LOGO_DIR.mkdir(parents=True, exist_ok=True)
    POI_IMAGE_DIR.mkdir(parents=True, exist_ok=True)

    if not POI_DB_PATH.exists():
        seed_path = BASE_DIR / "data" / "poi_database.csv"
        if seed_path.exists() and seed_path.resolve() != POI_DB_PATH.resolve():
            shutil.copyfile(seed_path, POI_DB_PATH)
        else:
            normalize_poi_df(pd.DataFrame(DEFAULT_POIS)).to_csv(POI_DB_PATH, index=False)

    if not POI_TEMPLATE_PATH.exists():
        template = pd.DataFrame([{column: "" for column in POI_COLUMNS}])
        template.loc[0, "poi_id"] = "POI-NEW-001"
        template.loc[0, "currency"] = "CNY"
        template.loc[0, "tags"] = "科技;文化;雨天可去"
        template.loc[0, "suitable_for"] = "商务客户;家庭客户"
        template.loc[0, "image_path"] = "assets/poi_images/example.jpg"
        template.loc[0, "image_url"] = ""
        template.loc[0, "image_alt"] = "Example POI image description"
        template.loc[0, "image_source"] = ""
        template.loc[0, "image_notes"] = "适合 Day Detail 横版图"
        template.loc[0, "images"] = "[]"
        template.loc[0, "status"] = "active"
        normalize_poi_df(template).to_csv(POI_TEMPLATE_PATH, index=False)


def load_pois():
    """Load POIs from the local CSV database."""
    ensure_poi_files()
    return normalize_poi_df(pd.read_csv(POI_DB_PATH))


def save_pois(df):
    """Save POIs back to the local CSV database."""
    normalized = normalize_poi_df(df)
    normalized.to_csv(POI_DB_PATH, index=False)
    return normalized


def split_multi_values(value):
    """Split semicolon/comma separated tag fields for filtering."""
    text = clean_text(value).replace("，", ";").replace(",", ";")
    return [item.strip() for item in text.split(";") if item.strip()]


def unique_multi_values(df, column):
    values = []
    for value in df[column].tolist():
        values.extend(split_multi_values(value))
    return sorted(set(values))


def display_poi_df(df):
    """Show POI fields with Chinese column names in the Chinese interface."""
    display_df = df.copy()
    if "image_path" in display_df.columns:
        display_df["image_status"] = display_df.apply(lambda row: "有图片" if poi_has_image(row) else "暂无图片", axis=1)
    if "status" in display_df.columns:
        display_df["status"] = display_df["status"].apply(poi_status_label)
    display_df["operation"] = "编辑 / 详情管理 / 图片管理 / 删除"
    labels = {**POI_FIELD_LABELS, "image_status": "图片状态", "operation": "操作"}
    return display_df.rename(columns=labels)


def poi_column_config():
    """Keep English data keys internally while displaying Chinese labels."""
    return {
        "id": st.column_config.TextColumn("内部ID"),
        "poi_id": st.column_config.TextColumn("点位编号"),
        "code": st.column_config.TextColumn("点位编码"),
        "name_cn": st.column_config.TextColumn("中文名称"),
        "name_en": st.column_config.TextColumn("英文名称"),
        "category": st.column_config.TextColumn("点位类型"),
        "city": st.column_config.TextColumn("城市"),
        "district": st.column_config.TextColumn("区域"),
        "duration_hours": st.column_config.NumberColumn("建议游览时长（小时）", min_value=0.0, step=0.5),
        "price_min": st.column_config.NumberColumn("最低成本价", min_value=0.0, step=100.0),
        "price_max": st.column_config.NumberColumn("最高成本价", min_value=0.0, step=100.0),
        "currency": st.column_config.TextColumn("币种"),
        "tags": st.column_config.TextColumn("标签"),
        "suitable_for": st.column_config.TextColumn("适合人群"),
        "description": st.column_config.TextColumn("综合描述"),
        "description_cn": st.column_config.TextColumn("中文介绍"),
        "description_en": st.column_config.TextColumn("英文介绍"),
        "recommended_reason": st.column_config.TextColumn("推荐理由"),
        "highlights": st.column_config.TextColumn("亮点"),
        "target_users": st.column_config.TextColumn("目标客群"),
        "image_placeholder": st.column_config.TextColumn("图片占位符或图片路径"),
        "image_notes": st.column_config.TextColumn("图片备注"),
        "image_source": st.column_config.TextColumn("图片来源"),
        "image_alt": st.column_config.TextColumn("图片英文说明"),
        "image_url": st.column_config.TextColumn("网络图片地址"),
        "image_path": st.column_config.TextColumn("本地图片路径"),
        "images": st.column_config.TextColumn("图片组"),
        "address": st.column_config.TextColumn("地址"),
        "opening_hours": st.column_config.TextColumn("开放时间"),
        "reservation_info": st.column_config.TextColumn("预约信息"),
        "notes": st.column_config.TextColumn("内部备注"),
        "internal_remark": st.column_config.TextColumn("内部备注补充"),
        "status": st.column_config.SelectboxColumn("状态", options=["active", "deleted"]),
    }


def shorten_client_text(value, max_chars=150):
    """Shorten customer-facing text without exposing internal fields."""
    text = clean_text(value)
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rsplit(" ", 1)[0].rstrip(".,;") + "."


def client_safe_export_text(value):
    """Remove backend POI IDs, costs, durations, and placeholder words from PPT export copy."""
    import re

    text = clean_text(value)
    blocked_terms = [
        "Featured Program",
        "No POIs selected",
        "Selected POIs",
        "Image Placeholder",
        "Visual Highlight Image Placeholder",
        "placeholder",
        "Photo Area",
        "undefined",
        "null",
    ]
    for term in blocked_terms:
        text = re.sub(re.escape(term), "", text, flags=re.IGNORECASE)
    text = re.sub(r"\b[A-Za-z ]*POI[- ]?[A-Za-z0-9]+\b", "", text)
    text = re.sub(r"\bCNY\s*[0-9,]+(?:\s*-\s*CNY\s*[0-9,]+)?", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\b\d+(?:\.\d+)?h\b", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\b(price_min|price_max|duration_hours|poi_id|notes|category)\b", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text).strip(" -|,.;")
    return text


def output_poi_name(poi):
    """Use only the official English POI name in client-facing output."""
    return client_safe_export_text(poi.get("name_en"))


def output_poi_description(poi):
    """Use only the English POI description in client-facing output."""
    return client_safe_export_text(poi.get("description_en"))


def format_poi_for_client_export(poi):
    """Return only customer-visible POI fields for PPT export."""
    title = output_poi_name(poi)
    if not title:
        return None
    description = shorten_client_text(output_poi_description(poi), max_chars=145)
    return {"title": title, "description": description}


def is_english_safe_text(value):
    """Only use placeholder labels in the PPT when they are English-safe text."""
    text = clean_text(value)
    return bool(text) and all(ord(char) < 128 for char in text)


def filter_pois(df, categories, cities, tags, suitable_for, keyword):
    """Apply POI list filters used by the management page."""
    filtered = df.copy()
    if categories:
        filtered = filtered[filtered["category"].isin(categories)]
    if cities:
        filtered = filtered[filtered["city"].isin(cities)]
    if tags:
        filtered = filtered[filtered["tags"].apply(lambda value: any(tag in split_multi_values(value) for tag in tags))]
    if suitable_for:
        filtered = filtered[filtered["suitable_for"].apply(lambda value: any(item in split_multi_values(value) for item in suitable_for))]
    if keyword:
        keyword_lower = keyword.lower()
        filtered = filtered[
            filtered["name_cn"].str.lower().str.contains(keyword_lower, na=False)
            | filtered["name_en"].str.lower().str.contains(keyword_lower, na=False)
        ]
    return filtered


def poi_display_label(row):
    """Build a readable option label for POI selectors."""
    name = clean_text(row.get("name_en")) or clean_text(row.get("name_cn")) or clean_text(row.get("poi_id"))
    city = clean_text(row.get("city"))
    category = clean_text(row.get("category"))
    return f"{name} | {city} | {category} | {row.get('poi_id')}"


def selected_poi_records(poi_df, poi_ids):
    """Return selected POI rows as dictionaries, preserving selected order."""
    records = []
    for poi_id in poi_ids or []:
        match = poi_df[poi_df["poi_id"].astype(str) == str(poi_id)]
        if not match.empty:
            records.append(match.iloc[0].to_dict())
    return records


def suggested_pois_for_day(city, theme, requirements=None, count=2):
    """Create clearly marked AI suggestions when the local POI database has no match."""
    city = clean_text(city) or "Selected City"
    theme_text = clean_text(theme)
    candidates = []
    if is_medical_theme(theme_text):
        candidates = [
            f"{city} International Health Checkup Center",
            f"{city} Recovery-friendly Wellness Lounge",
        ]
    elif keyword_hit(theme_text, ["Technology", "Business", "Innovation", "科技", "商务"]):
        candidates = [
            f"{city} Technology Enterprise Visit",
            f"{city} Innovation District Experience",
        ]
    elif keyword_hit(theme_text, ["Culture", "Cultural", "Heritage", "文化"]):
        candidates = [
            f"{city} Cultural Landmark Experience",
            f"{city} Local Lifestyle Visit",
        ]
    elif keyword_hit(theme_text, ["Wellness", "Recovery", "康养"]):
        candidates = [
            f"{city} Wellness Experience",
            f"{city} Healthy Dining Arrangement",
        ]
    elif keyword_hit(theme_text, ["Arrival", "Departure"]):
        candidates = [
            f"{city} Private Transfer Service",
            f"{city} Hotel Check-in Support",
        ]
    else:
        candidates = [
            f"{city} Curated Local Experience",
            f"{city} Concierge Dining Arrangement",
        ]
    return [
        {"name": name, "source": "ai_suggestion", "status": "未入库", "city": city}
        for name in candidates[:count]
    ]


def ensure_proposal_suggested_pois(proposal):
    """Attach non-database AI suggestion labels for days without selected database POIs."""
    for day in proposal.get("itinerary") or []:
        if day.get("poi_ids"):
            day.setdefault("suggested_pois", [])
        elif not day.get("suggested_pois"):
            day["suggested_pois"] = suggested_pois_for_day(day.get("city"), day.get("theme"))
    return proposal


def poi_source_badges_html(db_pois=None, suggested_pois=None):
    """Render source-colored POI badges for the internal editor."""
    badges = []
    for poi in db_pois or []:
        label = clean_text(poi.get("name_cn")) or clean_text(poi.get("name_en")) or clean_text(poi.get("poi_id"))
        if label:
            badges.append(f'<span class="poi-source-badge db">已入库 · {html.escape(label)}</span>')
    for poi in suggested_pois or []:
        if isinstance(poi, str):
            label = clean_text(poi)
        else:
            label = clean_text(poi.get("name")) or clean_text(poi.get("name_en")) or clean_text(poi.get("title"))
        if label:
            badges.append(f'<span class="poi-source-badge ai">AI建议 · 未入库 · {html.escape(label)}</span>')
    if not badges:
        return ""
    return '<div class="poi-source-badges">' + "".join(badges) + "</div>"


def poi_price_range(poi):
    """Format a POI cost range for display and PPT output."""
    currency = clean_text(poi.get("currency")) or "CNY"
    price_min = optional_number(poi.get("price_min"))
    price_max = optional_number(poi.get("price_max"))
    if price_min is not None and price_max is not None and price_min != price_max:
        return f"{money(price_min, currency)} - {money(price_max, currency)}"
    if price_min is not None or price_max is not None:
        return money(price_min if price_min is not None else price_max, currency)
    return "Cost TBD"


def poi_cost_bounds(day_pois):
    """Calculate min/max POI cost subtotal for later route costing."""
    min_total = 0
    max_total = 0
    for pois in day_pois or []:
        for poi in pois:
            price_min = optional_number(poi.get("price_min"))
            price_max = optional_number(poi.get("price_max"))
            min_total += price_min if price_min is not None else (price_max or 0)
            max_total += price_max if price_max is not None else (price_min or 0)
    return min_total, max_total


def format_cost_range(min_total, max_total, currency="CNY"):
    if min_total and max_total and min_total != max_total:
        return f"{money(min_total, currency)} - {money(max_total, currency)}"
    if min_total or max_total:
        return money(min_total or max_total, currency)
    return "To be calculated"


def parse_route_cities(route_cities):
    """Parse route city text into title-cased English city names."""
    raw_text = clean_text(route_cities)
    if not raw_text:
        return []
    for separator in ["，", "、", ";", "；", "/", "|"]:
        raw_text = raw_text.replace(separator, ",")
    cities = []
    for city in raw_text.split(","):
        normalized = " ".join(part.capitalize() for part in clean_text(city).split())
        if normalized and normalized not in cities:
            cities.append(normalized)
    return cities


def city_for_day(day_index, total_days, cities):
    """Distribute route cities across itinerary days without inventing cities."""
    if not cities:
        return ""
    if total_days <= len(cities):
        return cities[day_index]
    city_index = int(day_index * len(cities) / total_days)
    return cities[min(city_index, len(cities) - 1)]


def generated_theme(day_number, total_days, city):
    """Create an English placeholder theme for the generated itinerary draft."""
    if day_number == 1:
        return "Arrival & Welcome"
    if day_number == total_days:
        return "Departure Preparation"
    city_lower = city.lower()
    if "hangzhou" in city_lower:
        return "Medical Consultation"
    if "yiwu" in city_lower:
        return "Business & Local Visit"
    if "shanghai" in city_lower:
        return "City Experience & Recovery"
    if "guangzhou" in city_lower:
        return "Business & Innovation Visit"
    return "Curated Local Experience"


def generated_arrangement(day_number, total_days, city):
    """Create English placeholder arrangement copy for generated itinerary rows."""
    if day_number == 1:
        return "Airport reception, private transfer, hotel check-in, and welcome briefing."
    if day_number == total_days:
        return "Final consultation, leisure time, and departure preparation."
    city_lower = city.lower()
    if "hangzhou" in city_lower:
        return "Coordinated medical consultation, translation support, and light recovery-friendly itinerary."
    if "yiwu" in city_lower:
        return "Curated business visit, local market experience, and private transfer support."
    if "shanghai" in city_lower:
        return "Curated city visit, premium local experience, and recovery-friendly leisure time."
    return "Curated city visit and local experience based on client needs."



def english_list(items, limit=3):
    """Join selected English-safe names into proposal copy."""
    safe_items = [clean_text(item) for item in items if clean_text(item)]
    if not safe_items:
        return ""
    safe_items = safe_items[:limit]
    if len(safe_items) == 1:
        return safe_items[0]
    return ", ".join(safe_items[:-1]) + f" and {safe_items[-1]}"


def normalize_match_text(*values):
    """Build a lowercase search string for simple POI matching."""
    return " ".join(clean_text(value).lower() for value in values if clean_text(value))


def keyword_hit(text, keywords):
    """Check whether any Chinese or English keyword appears in a POI text bundle."""
    lowered = clean_text(text).lower()
    return any(clean_text(keyword).lower() in lowered for keyword in keywords if clean_text(keyword))


def route_day_profile(day_number, total_days, city, requirements):
    """Decide the kind of day to generate before selecting POIs."""
    customer_type = clean_text(requirements.get("customerType"))
    preference_tags = requirements.get("preferenceTags") or []
    required_categories = requirements.get("requiredCategories") or []
    city_lower = clean_text(city).lower()

    if day_number == 1:
        return {
            "theme": "Arrival & Welcome",
            "categories": ["交通", "接机", "酒店", "住宿", "餐厅", "餐费", "服务标准"],
            "tags": ["轻体力", "室内", "高端"],
            "max_pois": 2,
            "max_hours": 4,
        }
    if day_number == total_days:
        return {
            "theme": "Departure Preparation",
            "categories": ["交通", "送机", "酒店", "住宿", "服务标准", "餐费"],
            "tags": ["轻体力", "室内"],
            "max_pois": 1,
            "max_hours": 3,
        }
    medical_day = 2 if total_days > 2 else 1
    needs_medical = "医疗体检" in preference_tags or customer_type in ["康养客户", "综合客户"]
    if needs_medical and day_number == medical_day:
        return {
            "theme": "Group Health Checkup & Recovery",
            "categories": ["医疗机构", "体检", "康养体验", "餐厅", "餐费", "交通", "服务标准"],
            "tags": ["医疗", "体检", "康养", "轻体力", "室内", "健康餐"],
            "max_pois": 2,
            "max_hours": 6,
        }
    if customer_type in ["商务客户", "科技考察客户"] or "商务" in preference_tags or "科技" in preference_tags or "yiwu" in city_lower:
        return {
            "theme": "Business & Innovation Visit",
            "categories": ["科技企业", "商务参访", "产业园区", "餐厅", "餐费", "文化景点", "旅游", "服务标准"],
            "tags": ["商务", "科技", "名企", "产业考察", "高端"],
            "max_pois": 3,
            "max_hours": 8,
        }
    if customer_type == "康养客户" or "康养" in preference_tags:
        return {
            "theme": "Wellness & Light Recovery",
            "categories": ["康养体验", "餐厅", "餐费", "文化景点", "旅游", "酒店", "住宿"],
            "tags": ["康养", "轻体力", "健康餐", "茶文化", "室内"],
            "max_pois": 2,
            "max_hours": 6,
        }
    return {
        "theme": "Cultural & Local Experience",
        "categories": ["文化景点", "博物馆", "非遗体验", "夜游", "餐厅", "餐费", "旅游"],
        "tags": ["文化", "轻体力", "夜游", "美食", "适合外宾"],
        "max_pois": 3,
        "max_hours": 8,
    }


def generic_client_arrangement(day_number, total_days, city, theme):
    """Create customer-facing day copy without internal POI identifiers."""
    if day_number == 1:
        return f"Arrive in {city} with coordinated reception, private transfer, hotel check-in, and a welcome briefing."
    if day_number == total_days:
        return f"Prepare for departure with light arrangements, flexible leisure time, and private transfer support."
    if any(word in theme for word in ["Medical", "Consultation", "Health Checkup", "Screening"]):
        return f"Coordinated health checkups for all travelers in {city}, with bilingual support, private transfers, and recovery-friendly arrangements."
    if "Business" in theme or "Innovation" in theme:
        return f"Explore {city}'s business and innovation landscape through curated visits and local host support."
    if "Wellness" in theme or "Recovery" in theme:
        return f"Enjoy a recovery-friendly wellness day in {city}, featuring light cultural experiences and relaxing local activities."
    if "Cultural" in theme or "Culture" in theme:
        return f"Explore {city}'s local culture and lifestyle through a curated, comfortable day program."
    return generated_arrangement(day_number, total_days, city)


def build_arrangement_from_pois(day_number, total_days, city, theme, selected_pois):
    """Write English arrangement copy using only official English POI names when available."""
    client_pois = [format_poi_for_client_export(poi) for poi in selected_pois or []]
    poi_names = [poi["title"] for poi in client_pois if poi]
    name_text = english_list(poi_names, limit=2)
    if selected_pois and name_text:
        if day_number == 1:
            return f"Arrive in {city} with coordinated reception and a gentle welcome program featuring {name_text}."
        if day_number == total_days:
            return f"Prepare for departure with light arrangements around {name_text}, keeping the schedule relaxed and flexible."
        if any(word in theme for word in ["Medical", "Consultation", "Health Checkup", "Screening"]):
            return f"Coordinate group health checkups for all travelers in {city}, including {name_text}, with bilingual support, private transfers, and recovery-friendly arrangements."
        if "Business" in theme:
            return f"Explore {city}'s business and innovation landscape through curated visits including {name_text}."
        if "Wellness" in theme:
            return f"Enjoy a recovery-friendly wellness day in {city}, featuring {name_text} and relaxing local activities."
        return f"Explore {city}'s local culture and lifestyle through curated visits to {name_text}."
    return generic_client_arrangement(day_number, total_days, city, theme)


def score_poi_for_day(poi, city, profile, requirements, used_ids):
    """Score one POI for one generated route day."""
    poi_id = clean_text(poi.get("poi_id"))
    category = clean_text(poi.get("category"))
    poi_city = clean_text(poi.get("city"))
    tags = split_multi_values(poi.get("tags"))
    suitable_for = split_multi_values(poi.get("suitable_for"))
    text_bundle = normalize_match_text(
        poi.get("name_cn"), poi.get("name_en"), category, poi_city,
        poi.get("district"), poi.get("tags"), poi.get("suitable_for"),
        poi.get("description_cn"), poi.get("description_en"), poi.get("notes"),
    )

    excluded = requirements.get("excludedCategories") or []
    if any(clean_text(item) and (item == category or keyword_hit(text_bundle, [item])) for item in excluded):
        return -999
    if not poi_allowed_for_city(poi, city):
        return -999

    score = 60

    profile_category_match = keyword_hit(category, profile.get("categories", [])) or keyword_hit(text_bundle, profile.get("categories", []))
    if profile_category_match:
        score += 32
    else:
        score -= 18
    for tag in profile.get("tags", []):
        if tag in tags or keyword_hit(text_bundle, [tag]):
            score += 10

    customer_type = clean_text(requirements.get("customerType"))
    if customer_type and (customer_type in suitable_for or keyword_hit(text_bundle, [customer_type])):
        score += 18

    for tag in requirements.get("preferenceTags") or []:
        if tag in tags or keyword_hit(text_bundle, [tag]):
            score += 12

    for required in requirements.get("requiredCategories") or []:
        if required == category or keyword_hit(text_bundle, [required]):
            score += 12

    budget_note = clean_text(requirements.get("budgetNote"))
    budget_max = parse_optional_amount(requirements.get("budgetMax"))
    price_max = optional_number(poi.get("price_max")) or optional_number(poi.get("price_min"))
    if keyword_hit(budget_note, ["高", "高端", "品质", "premium", "luxury"]):
        if keyword_hit(text_bundle, ["高端", "精品", "私密", "深度"]):
            score += 10
    if budget_max is not None and price_max is not None:
        per_day_reference = max(budget_max / max(int(requirements.get("days") or 1), 1), 1)
        if price_max > per_day_reference:
            score -= 8

    if poi_id in used_ids:
        score -= 35
    return score


def select_pois_for_day(poi_records, city, profile, requirements, used_ids):
    """Choose a small, editable set of POIs for a generated day."""
    scored = []
    for poi in poi_records:
        score = score_poi_for_day(poi, city, profile, requirements, used_ids)
        if score > 0:
            scored.append((score, poi))
    scored.sort(key=lambda item: item[0], reverse=True)

    selected = []
    total_hours = 0.0
    max_pois = profile.get("max_pois", 3)
    max_hours = profile.get("max_hours", 8)
    # First pass prefers the itinerary city strongly.
    for _, poi in scored:
        if len(selected) >= max_pois:
            break
        if clean_text(poi.get("poi_id")) in used_ids:
            continue
        if not city_matches(poi.get("city"), city):
            continue
        duration = optional_number(poi.get("duration_hours")) or 1.0
        if total_hours + duration <= max_hours or not selected:
            selected.append(poi)
            total_hours += duration
            used_ids.add(clean_text(poi.get("poi_id")))

    return selected


def calculate_route_pricing_preview(day_pois):
    """Create an initial pricing preview and warnings from auto-selected POIs."""
    missing_price_count = 0
    for pois in day_pois or []:
        for poi in pois:
            if optional_number(poi.get("price_min")) is None and optional_number(poi.get("price_max")) is None:
                missing_price_count += 1
    result = calculate_pricing_preview(day_pois, None, None, None, None, None, "")
    result["missing_poi_price_count"] = missing_price_count
    result["warnings"] = []
    if missing_price_count:
        result["warnings"].append("部分点位缺少成本价格，当前报价仅供内部参考。")
    if not result.get("is_complete"):
        result["warnings"].append("尚未生成完整报价，请在确认具体线路和点位后进行价格核算。")
    return result


def generateRouteTemplateFromRequirements(requirements, poiDatabase):
    """Generate an editable route draft and automatic POI selections from client requirements."""
    total_days = int(requirements.get("days") or 1)
    cities = requirements.get("cities") or parse_route_cities(requirements.get("routeCities"))
    poi_records = poiDatabase.to_dict("records") if hasattr(poiDatabase, "to_dict") else (poiDatabase or [])
    used_ids = set()
    rows = []
    day_poi_ids = []
    day_pois = []

    for index in range(total_days):
        day_number = index + 1
        city = city_for_day(index, total_days, cities)
        profile = route_day_profile(day_number, total_days, city, requirements)
        intensity = clean_text(requirements.get("intensityLevel"))
        if intensity == "轻松":
            profile["max_pois"] = max(1, min(profile.get("max_pois", 3), 2))
            profile["max_hours"] = min(profile.get("max_hours", 8), 6)
        elif intensity == "紧凑" and day_number not in [1, total_days]:
            profile["max_pois"] = min(profile.get("max_pois", 3) + 1, 4)
            profile["max_hours"] = profile.get("max_hours", 8) + 1
        selected = select_pois_for_day(poi_records, city, profile, requirements, used_ids)
        theme = profile.get("theme") or generated_theme(day_number, total_days, city)
        arrangement = build_arrangement_from_pois(day_number, total_days, city, theme, selected)
        rows.append({"Day": day_number, "City": city, "Theme": theme, "Arrangement": arrangement})
        day_pois.append(selected)
        day_poi_ids.append([clean_text(poi.get("poi_id")) for poi in selected if clean_text(poi.get("poi_id"))])

    rows, day_poi_ids = validate_generated_route(rows, day_poi_ids, poiDatabase, requirements)
    day_pois = [selected_poi_records(poiDatabase, ids) for ids in day_poi_ids] if hasattr(poiDatabase, "iterrows") else day_pois
    return {
        "days": rows,
        "dayPoiIds": day_poi_ids,
        "dayPois": day_pois,
        "pricingPreview": calculate_route_pricing_preview(day_pois),
    }



def deepseek_api_key():
    """Read DeepSeek API key from Streamlit secrets or environment variables."""
    try:
        key = st.secrets.get("DEEPSEEK_API_KEY", "")
        if key:
            key = str(key).strip()
        else:
            key = ""
    except Exception:
        key = ""
    if not key:
        key = os.getenv("DEEPSEEK_API_KEY", "").strip()
    placeholder_words = ["你的", "your", "api key", "sk-xxxx", "填入"]
    if any(word in key.lower() for word in placeholder_words):
        return ""
    try:
        key.encode("ascii")
    except UnicodeEncodeError:
        return ""
    return key


def call_deepseek_json(system_prompt, user_prompt, fallback=None):
    """Call DeepSeek through its OpenAI-compatible chat endpoint and parse JSON output."""
    key = deepseek_api_key()
    if not key:
        return fallback
    payload = {
        "model": os.getenv("DEEPSEEK_MODEL", "deepseek-chat"),
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "temperature": 0.35,
        "response_format": {"type": "json_object"},
    }
    request = urllib.request.Request(
        "https://api.deepseek.com/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=45) as response:
            data = json.loads(response.read().decode("utf-8"))
        content = data.get("choices", [{}])[0].get("message", {}).get("content", "{}")
        return json.loads(content)
    except (urllib.error.URLError, urllib.error.HTTPError, TimeoutError, UnicodeEncodeError, json.JSONDecodeError, KeyError) as exc:
        st.warning(f"DeepSeek 调用失败，已使用本地规则生成：{exc}")
        return fallback


def ordered_city_mentions(text_value):
    """Return route cities in the order they appear in the user's natural-language request."""
    source = clean_text(text_value)
    lower_source = source.lower()
    city_candidates = [
        ("Shanghai", ["上海", "shanghai"]),
        ("Hangzhou", ["杭州", "hangzhou"]),
        ("Yiwu", ["义乌", "yiwu"]),
        ("Guangzhou", ["广州", "guangzhou"]),
    ]
    hits = []
    for city, keys in city_candidates:
        positions = []
        for key in keys:
            haystack = lower_source if key.isascii() else source
            needle = key.lower() if key.isascii() else key
            position = haystack.find(needle)
            if position >= 0:
                positions.append(position)
        if positions:
            hits.append((min(positions), city))
    return [city for _, city in sorted(hits, key=lambda item: item[0])]


def move_city_block_to_front(proposal, city_name):
    """Move all days for a city to the front while preserving each city's internal order."""
    target_key = city_key(city_name)
    if not target_key:
        return proposal
    itinerary = proposal.get("itinerary") or []
    front = [day for day in itinerary if city_key(day.get("city")) == target_key]
    rest = [day for day in itinerary if city_key(day.get("city")) != target_key]
    if not front:
        return proposal
    proposal["itinerary"] = front + rest
    proposal["selected_pois"] = [day.get("poi_ids") or [] for day in proposal["itinerary"]]
    return proposal


def requested_day_count(text_value):
    """Extract intended total trip days, prioritizing explicit trip-duration wording."""
    text_value = clean_text(text_value)
    patterns = [
        r"(?:一共|总共|总计|共|整体|全程|还是|保持|控制在|不要超过)\s*(\d+)\s*天",
        r"(\d+)\s*天\s*(?:行程|旅程|线路|旅行|游玩)",
        r"(?:行程|旅程|线路|旅行|游玩)\s*(?:一共|总共|共|是|为)?\s*(\d+)\s*天",
        r"(\d+)\s*(?:days?|day)",
        r"(\d+)\s*天",
    ]
    for pattern in patterns:
        match = re.search(pattern, text_value, flags=re.I)
        if match:
            try:
                return int(match.group(1))
            except (TypeError, ValueError):
                return None
    return None


def enforce_proposal_day_count(proposal, target_days, protected_cities=None, poi_database=None, requirements=None):
    """Resize itinerary to the requested total day count while protecting newly requested cities and medical days."""
    if not target_days or target_days <= 0:
        return proposal
    itinerary = proposal.get("itinerary") or []
    protected_keys = {city_key(city) for city in (protected_cities or []) if city_key(city)}

    def is_protected(day):
        return city_key(day.get("city")) in protected_keys or is_medical_theme(day.get("theme"))

    while len(itinerary) > target_days:
        removable_index = None
        for index in range(len(itinerary) - 1, -1, -1):
            if not is_protected(itinerary[index]):
                removable_index = index
                break
        if removable_index is None:
            removable_index = len(itinerary) - 1
        itinerary.pop(removable_index)

    while len(itinerary) < target_days:
        city = clean_text(itinerary[-1].get("city")) if itinerary else "Hangzhou"
        itinerary.append(create_city_itinerary_day(city, itinerary, requirements or {}, poi_database))

    for index, day in enumerate(itinerary, start=1):
        day["day"] = index
    proposal["itinerary"] = itinerary
    proposal["selected_pois"] = [day.get("poi_ids") or [] for day in itinerary]
    return proposal


def create_city_itinerary_day(city_name, itinerary, requirements, poi_database):
    """Create one itinerary day for an added city during conversational revision."""
    city = clean_text(city_name)
    day_number = len(itinerary) + 1
    total_days = day_number
    profile = route_day_profile(day_number, total_days, city, requirements or {})
    if profile.get("theme") in ["Departure Preparation", "Arrival & Welcome"]:
        profile["theme"] = generated_theme(max(2, day_number - 1), max(3, total_days + 1), city)
    poi_records = poi_database.to_dict("records") if hasattr(poi_database, "to_dict") else (poi_database or [])
    used_ids = {clean_text(poi_id) for day in itinerary for poi_id in (day.get("poi_ids") or []) if clean_text(poi_id)}
    selected = select_pois_for_day(poi_records, city, profile, requirements or {}, used_ids)
    theme = profile.get("theme") or generated_theme(max(2, day_number), max(3, total_days + 1), city)
    arrangement = build_arrangement_from_pois(day_number, total_days, city, theme, selected)
    poi_ids = [clean_text(poi.get("poi_id")) for poi in selected if clean_text(poi.get("poi_id"))]
    return {
        "day": day_number,
        "city": city,
        "theme": theme,
        "description": arrangement,
        "poi_ids": poi_ids,
        "suggested_pois": [] if poi_ids else suggested_pois_for_day(city, theme),
        "medical": "group" if is_medical_theme(theme) else "",
        "notes": "AI added during conversational revision",
    }


def add_city_day_to_proposal(proposal, city_name, requirements, poi_database):
    """Add a new city day before departure when the user asks to add a destination."""
    itinerary = proposal.get("itinerary") or []
    city = clean_text(city_name)
    new_day = create_city_itinerary_day(city, itinerary, requirements, poi_database)
    insert_index = len(itinerary)
    if itinerary and keyword_hit(clean_text(itinerary[-1].get("theme")), ["Departure", "离境", "返程"]):
        insert_index = max(len(itinerary) - 1, 0)
    itinerary.insert(insert_index, new_day)
    proposal["itinerary"] = itinerary
    proposal["selected_pois"] = [day.get("poi_ids") or [] for day in itinerary]
    return proposal


def infer_requirements_from_text(user_requirements, filters):
    """Extract basic structured requirements from natural language plus sidebar filters."""
    text_value = clean_text(user_requirements)
    lower = text_value.lower()
    inferred = dict(filters or {})
    inferred["raw_user_input"] = text_value

    extracted_days = requested_day_count(text_value)
    if extracted_days:
        inferred["days"] = extracted_days
        inferred["nights"] = max(extracted_days - 1, 0)
    people_match = re.search(r"(\d+)\s*(?:个人|人|位|traveler|travelers|guests?)", text_value, flags=re.I)
    if people_match:
        inferred["customerCount"] = int(people_match.group(1))

    mentioned_cities = ordered_city_mentions(text_value)
    if mentioned_cities:
        inferred["cities"] = mentioned_cities
        inferred["routeCities"] = ", ".join(mentioned_cities)

    preference_tags = list(inferred.get("preferenceTags") or [])
    for keyword, tag in [
        ("体检", "医疗体检"), ("医疗", "医疗体检"), ("康养", "康养"),
        ("科技", "科技"), ("商务", "商务"), ("文化", "文化"),
        ("西湖", "文化"), ("轻松", "轻体力"), ("不要太累", "轻体力"),
        ("高端", "高端"), ("亲子", "亲子"),
    ]:
        if keyword in text_value and tag not in preference_tags:
            preference_tags.append(tag)
    inferred["preferenceTags"] = preference_tags

    if "科技" in text_value and "商务" not in text_value:
        inferred["customerType"] = "科技考察客户"
    if "康养" in text_value or "不要太累" in text_value or "轻松" in text_value:
        inferred["intensityLevel"] = "轻松"
    if "家庭" in text_value or "亲子" in text_value:
        inferred["customerType"] = "家庭客户"
    if "商务" in text_value:
        inferred["customerType"] = "商务客户"

    if "科技企业" in text_value:
        required = list(inferred.get("requiredCategories") or [])
        if "科技企业" not in required:
            required.append("科技企业")
        inferred["requiredCategories"] = required
    return inferred


def requirement_summary(requirements):
    """Build a Chinese summary for the optional constraint sidebar."""
    req = requirements or {}
    rows = []
    mapping = [
        ("raw_user_input", "原始需求"),
        ("customerCount", "出行人数"),
        ("days", "出行天数"),
        ("routeCities", "城市范围"),
        ("customerType", "客户类型"),
        ("intensityLevel", "行程强度"),
        ("preferenceTags", "偏好标签"),
        ("requiredCategories", "必须包含"),
        ("excludedCategories", "不希望包含"),
        ("budgetNote", "预算说明"),
    ]
    for key, label in mapping:
        value = req.get(key)
        if isinstance(value, list):
            value = "、".join(clean_text(item) for item in value if clean_text(item))
        value = clean_text(value)
        if value:
            rows.append(f"- {label}：{value}")
    return "\n".join(rows) or "尚未识别。请先在右侧输入客户自然语言需求。"


def proposal_from_route_template(requirements, route_template, intro_text):
    """Create the unified current_proposal object from generated route data."""
    itinerary = []
    day_ids = route_template.get("dayPoiIds") or []
    for index, row in enumerate(route_template.get("days") or [], start=1):
        current_ids = day_ids[index - 1] if index - 1 < len(day_ids) else []
        city = clean_text(row.get("City"))
        theme = clean_text(row.get("Theme"))
        itinerary.append({
            "day": index,
            "city": city,
            "theme": theme,
            "description": clean_text(row.get("Arrangement")),
            "poi_ids": current_ids,
            "suggested_pois": [] if current_ids else suggested_pois_for_day(city, theme),
            "medical": "group" if is_medical_theme(row.get("Theme")) else "",
            "notes": "",
        })
    return {
        "proposal_title": clean_text(requirements.get("title")) or "Premium Medical Tourism Proposal",
        "proposal_intro": clean_text(intro_text) or "A premium health travel proposal with coordinated medical support and curated local experiences.",
        "client_profile": {
            "customer_count": int(requirements.get("customerCount") or 1),
            "customer_type": clean_text(requirements.get("customerType")),
            "requirements_text": clean_text(requirements.get("requirementsText")),
            "intensity_level": clean_text(requirements.get("intensityLevel")),
        },
        "itinerary": itinerary,
        "medical_plan": "Group health checkups for all travelers should be arranged on the same medical day unless special requirements require otherwise.",
        "selected_pois": day_ids,
        "service_inclusions": "Private airport transfers\nMedical appointment coordination\nLocal host and translation support\nCurated sightseeing arrangements",
        "cost_notes": clean_text(requirements.get("clientQuote")) or "To be confirmed",
        "next_steps": "Confirm travel dates\nConfirm health checkup package\nConfirm hotel preference\nFinalize quotation",
        "export_settings": {"output_language": "en"},
        "locked_modules": [],
    }


def proposal_to_session_state(proposal, pricing_preview=None):
    """Sync current_proposal into existing Streamlit state used by editors and exports."""
    proposal = ensure_proposal_suggested_pois(proposal)
    itinerary = proposal.get("itinerary") or []
    rows = []
    day_ids = []
    day_suggested = []
    for index, day in enumerate(itinerary, start=1):
        rows.append({
            "Day": index,
            "City": clean_text(day.get("city")),
            "Theme": clean_text(day.get("theme")),
            "Arrangement": clean_text(day.get("description")),
        })
        day_ids.append(day.get("poi_ids") or [])
        day_suggested.append(day.get("suggested_pois") or [])
    st.session_state.current_proposal = proposal
    st.session_state.generated_intro = clean_text(proposal.get("proposal_intro"))
    st.session_state.itinerary_rows = rows
    st.session_state.day_poi_ids = day_ids
    st.session_state.day_suggested_pois = day_suggested
    if pricing_preview is not None:
        st.session_state.route_pricing_preview = pricing_preview
    st.session_state.generationStatus = "completed"
    st.session_state.generation_error = ""
    st.session_state.draft_version += 1


def sync_current_proposal_from_session():
    """Keep current_proposal aligned with manual edits in the structured editor."""
    proposal = st.session_state.get("current_proposal") or {}
    itinerary = []
    for index, row in enumerate(itinerary_records(st.session_state.get("itinerary_rows", [])), start=1):
        ids = st.session_state.day_poi_ids[index - 1] if index - 1 < len(st.session_state.get("day_poi_ids", [])) else []
        suggested = st.session_state.day_suggested_pois[index - 1] if index - 1 < len(st.session_state.get("day_suggested_pois", [])) else []
        itinerary.append({
            "day": index,
            "city": clean_text(row.get("City")),
            "theme": clean_text(row.get("Theme")),
            "description": clean_text(row.get("Arrangement")),
            "poi_ids": ids,
            "suggested_pois": [] if ids else suggested_pois_for_day(row.get("City"), row.get("Theme")) if not suggested else suggested,
            "medical": "group" if is_medical_theme(row.get("Theme")) else "",
            "notes": "",
        })
    proposal["proposal_intro"] = st.session_state.get("generated_intro", proposal.get("proposal_intro", ""))
    proposal["itinerary"] = itinerary
    proposal["selected_pois"] = st.session_state.get("day_poi_ids", [])
    st.session_state.current_proposal = proposal
    return proposal


def append_proposal_version(label):
    """Save a lightweight version snapshot."""
    history = st.session_state.setdefault("proposal_versions", [])
    snapshot = json.loads(json.dumps(st.session_state.get("current_proposal") or {}, ensure_ascii=False))
    history.append({"label": label, "proposal": snapshot, "time": time.strftime("%H:%M:%S")})


def generate_initial_proposal(user_requirements, filters, poi_database):
    """Generate first proposal version from natural language and sidebar filters."""
    inferred = infer_requirements_from_text(user_requirements, filters)
    inferred["requirementsText"] = user_requirements
    llm_result = call_deepseek_json(
        "You are a medical tourism proposal planner. Return JSON only with optional keys: proposal_title, proposal_intro, customer_type, intensity_level.",
        json.dumps({"user_requirements": user_requirements, "filters": inferred}, ensure_ascii=False),
        fallback={},
    ) or {}
    if clean_text(llm_result.get("proposal_title")):
        inferred["title"] = clean_text(llm_result.get("proposal_title"))
    if clean_text(llm_result.get("customer_type")):
        inferred["customerType"] = clean_text(llm_result.get("customer_type"))
    if clean_text(llm_result.get("intensity_level")):
        inferred["intensityLevel"] = clean_text(llm_result.get("intensity_level"))
    intro_text = clean_text(llm_result.get("proposal_intro")) or clean_text(filters.get("intro"))
    route_template = generateRouteTemplateFromRequirements(inferred, poi_database)
    proposal = proposal_from_route_template(inferred, route_template, intro_text)
    return proposal, route_template


def revise_proposal(current_proposal, user_instruction, filters, poi_database):
    """Revise current proposal in place according to a follow-up natural-language instruction."""
    proposal = json.loads(json.dumps(current_proposal or {}, ensure_ascii=False))
    instruction = clean_text(user_instruction)
    target_days = requested_day_count(instruction) or int((filters or {}).get("days") or len(proposal.get("itinerary") or []) or 0)
    itinerary = proposal.get("itinerary") or []
    locked = set(proposal.get("locked_modules") or [])

    if "proposal_intro" not in locked and any(word in instruction for word in ["简介", "高级", "高端", "更好", "优化英文"]):
        fallback_intro = "This premium health travel proposal combines coordinated medical checkups, attentive bilingual support, private transfers, and curated city experiences designed for a smooth and confidence-building journey."
        llm_result = call_deepseek_json(
            "Return JSON only with key proposal_intro. Write polished English client-facing copy, no Chinese.",
            json.dumps({"current_intro": proposal.get("proposal_intro"), "instruction": instruction}, ensure_ascii=False),
            fallback={"proposal_intro": fallback_intro},
        ) or {}
        proposal["proposal_intro"] = clean_text(llm_result.get("proposal_intro")) or fallback_intro

    if "itinerary" not in locked:
        added_city_requests = ordered_city_mentions(instruction)
        if added_city_requests and any(word in instruction for word in ["增加", "加入", "添加", "安排", "加上", "去"]):
            for city_name in added_city_requests:
                proposal = add_city_day_to_proposal(proposal, city_name, filters, poi_database)
            itinerary = proposal.get("itinerary") or []
        if "义乌减少一天" in instruction or ("义乌" in instruction and "减少" in instruction):
            yiwu_indices = [i for i, day in enumerate(itinerary) if "yiwu" in clean_text(day.get("city")).lower()]
            if len(yiwu_indices) > 1:
                itinerary.pop(yiwu_indices[-1])
        if "西湖" in instruction:
            target = next((day for day in itinerary if "hangzhou" in clean_text(day.get("city")).lower() and not is_medical_theme(day.get("theme"))), None)
            if target is None:
                target = {"city": "Hangzhou", "theme": "West Lake Cultural Experience", "description": "Enjoy a relaxed West Lake cultural experience with light walking, scenic moments, and local concierge support.", "poi_ids": [], "medical": "", "notes": ""}
                itinerary.insert(min(2, len(itinerary)), target)
            target["theme"] = "West Lake Cultural Experience"
            target["description"] = "Enjoy a relaxed West Lake cultural experience with light walking, scenic views, and recovery-friendly local support."
        if any(word in instruction for word in ["不想太商务", "更轻松", "不要太累", "轻松一点"]):
            for day in itinerary:
                if "Business" in clean_text(day.get("theme")) and "科技" not in instruction:
                    day["theme"] = "Light Cultural & Wellness Experience"
                    day["description"] = f"Enjoy a relaxed day in {day.get('city')}, with light cultural experiences, comfortable pacing, and local concierge support."
        if "体检" in instruction and any(word in instruction for word in ["同一天", "一起", "并行", "两个人", "2个人"]):
            medical_days = [day for day in itinerary if is_medical_theme(day.get("theme"))]
            if not medical_days and itinerary:
                medical_days = [itinerary[min(1, len(itinerary) - 1)]]
            if medical_days:
                primary = medical_days[0]
                primary["theme"] = "Group Health Checkup & Recovery"
                primary["medical"] = "group"
                primary["description"] = f"Coordinated health checkups for all travelers in {primary.get('city')}, with bilingual support, private transfers, and recovery-friendly arrangements."
                for duplicate in medical_days[1:]:
                    duplicate["theme"] = "Wellness & Local Experience"
                    duplicate["medical"] = ""
                    duplicate["description"] = f"Enjoy a relaxed recovery-friendly day in {duplicate.get('city')}, with light local experiences and concierge support."
        day_match = re.search(r"day\s*(\d+)|第\s*(\d+)\s*天", instruction, flags=re.I)
        if day_match and any(word in instruction for word in ["重写", "修改", "调整", "rewrite"]):
            day_no = int(day_match.group(1) or day_match.group(2))
            if 1 <= day_no <= len(itinerary):
                day = itinerary[day_no - 1]
                day["description"] = f"A refined, client-friendly day in {day.get('city')} adjusted according to the latest consultant instruction: {instruction}"

    first_city_match = re.search(r"(?:先到|先去|先飞到|飞到|从|第一站|先).*?(上海|杭州|义乌|广州|Shanghai|Hangzhou|Yiwu|Guangzhou)", instruction, flags=re.I)
    if first_city_match:
        proposal["itinerary"] = itinerary
        proposal = move_city_block_to_front(proposal, first_city_match.group(1))
        itinerary = proposal.get("itinerary") or []

    proposal["itinerary"] = itinerary
    proposal = enforce_proposal_day_count(proposal, target_days, protected_cities=ordered_city_mentions(instruction), poi_database=poi_database, requirements=filters)
    itinerary = proposal.get("itinerary") or []
    rows = [{"Day": i + 1, "City": day.get("city"), "Theme": day.get("theme"), "Arrangement": day.get("description")} for i, day in enumerate(itinerary)]
    ids = [day.get("poi_ids") or [] for day in itinerary]
    rows, ids = validate_generated_route(rows, ids, poi_database, filters)
    existing_suggested = [day.get("suggested_pois") or [] for day in itinerary]
    proposal["itinerary"] = [
        {
            "day": row.get("Day"),
            "city": row.get("City"),
            "theme": row.get("Theme"),
            "description": row.get("Arrangement"),
            "poi_ids": ids[index] if index < len(ids) else [],
            "suggested_pois": [] if (index < len(ids) and ids[index]) else (existing_suggested[index] if index < len(existing_suggested) and existing_suggested[index] else suggested_pois_for_day(row.get("City"), row.get("Theme"))),
            "medical": "group" if is_medical_theme(row.get("Theme")) else "",
            "notes": "",
        }
        for index, row in enumerate(rows)
    ]
    proposal["selected_pois"] = ids
    return proposal


def generate_itinerary_from_basics(days, route_cities):
    """Generate editable itinerary draft rows from sidebar basics."""
    total_days = int(days or 1)
    cities = parse_route_cities(route_cities)
    rows = []
    for index in range(total_days):
        day_number = index + 1
        city = city_for_day(index, total_days, cities)
        rows.append(
            {
                "Day": day_number,
                "City": city,
                "Theme": generated_theme(day_number, total_days, city),
                "Arrangement": generated_arrangement(day_number, total_days, city),
            }
        )
    return rows


def proposal_basics_key(
    title,
    group_size,
    days,
    nights,
    route_cities,
    budget_min,
    budget_max,
    budget_notes,
    client_quote,
    intro,
    customer_type,
    preference_tags,
    intensity_level,
    required_categories,
    excluded_categories,
):
    """Track whether left-side demand and budget-reference inputs changed after draft generation."""
    return "|".join(
        [
            clean_text(title),
            clean_text(group_size),
            clean_text(days),
            clean_text(nights),
            clean_text(route_cities),
            clean_text(budget_min),
            clean_text(budget_max),
            clean_text(budget_notes),
            clean_text(client_quote),
            clean_text(intro),
            clean_text(customer_type),
            ";".join(preference_tags or []),
            clean_text(intensity_level),
            ";".join(required_categories or []),
            ";".join(excluded_categories or []),
        ]
    )


def parse_optional_amount(value):
    """Parse optional budget/cost fields without turning blanks into zero."""
    text = clean_text(value).replace(",", "")
    if not text:
        return None
    try:
        return float(text)
    except ValueError:
        return None


def format_optional_usd(value):
    amount = optional_number(value)
    if amount is None:
        return "To be confirmed"
    return money(amount, "USD")


def format_client_quote(value):
    """Display the single client-facing quote as free text, or a clean fallback."""
    text = clean_text(value)
    return text if text else "To be confirmed"


def calculate_pricing_preview(day_pois, hotel_cost, transport_cost, service_cost, optional_addon_cost, profit_multiplier, client_quote):
    """Calculate right-side pricing preview from selected itinerary content and one final client quote."""
    poi_min, poi_max = poi_cost_bounds(day_pois)
    hotel = parse_optional_amount(hotel_cost)
    transport = parse_optional_amount(transport_cost)
    service = parse_optional_amount(service_cost)
    optional_addon = parse_optional_amount(optional_addon_cost)
    multiplier = parse_optional_amount(profit_multiplier) or 1.3
    client_quote_text = clean_text(client_quote)

    known_components = [value for value in [hotel, transport, service, optional_addon] if value is not None]
    has_poi_cost = bool(poi_min or poi_max)
    has_any_cost = bool(known_components or has_poi_cost)
    base_extra = sum(known_components)
    total_min = (poi_min if has_poi_cost else 0) + base_extra if has_any_cost else None
    total_max = (poi_max if has_poi_cost else 0) + base_extra if has_any_cost else None
    suggested_price = total_max * multiplier if total_max is not None else None

    return {
        "base_cost": None,
        "poi_cost_min": poi_min,
        "poi_cost_max": poi_max,
        "hotel_cost": hotel,
        "transport_cost": transport,
        "service_cost": service,
        "optional_addon_cost": optional_addon,
        "total_cost_min": total_min,
        "total_cost_max": total_max,
        "profit_multiplier": multiplier,
        "suggested_price": suggested_price,
        "client_quote": client_quote_text,
        "is_complete": has_any_cost,
    }


def validate_proposal_basics(title, group_size, days, route_cities):
    """Validate required inputs before generating the editable draft."""
    errors = []
    if not clean_text(title):
        errors.append("请填写英文方案标题")
    if not group_size or int(group_size) <= 0:
        errors.append("请填写客户人数")
    if not days or int(days) <= 0:
        errors.append("请填写天数")
    if not parse_route_cities(route_cities):
        errors.append("请填写英文路线城市")
    return errors


def poi_summary_text(pois, max_items=3):
    """Create customer-friendly Highlights text for daily slides."""
    client_items = []
    for poi in pois or []:
        formatted = format_poi_for_client_export(poi)
        if formatted:
            client_items.append(formatted)
    if not client_items:
        return ""

    lines = []
    for item in client_items[:max_items]:
        line = f"- {item['title']}"
        if item.get("description"):
            line += f"\n  {item['description']}"
        lines.append(line)
    return "\n".join(lines)


def poi_image_label(day, pois):
    """Keep client PPT image areas clean without backend placeholder wording."""
    return ""


def image_bytes_from_poi_placeholder(pois):
    """Use the first selected POI image as a client-facing fallback image."""
    preferred = []
    regular = []
    for poi in pois or []:
        notes = clean_text(poi.get("image_notes")).lower()
        if "cover" in notes or "main" in notes or "主图" in notes:
            preferred.append(poi)
        else:
            regular.append(poi)
    for poi in preferred + regular:
        for field in ["image_path", "image_placeholder"]:
            image_path = resolve_local_image_path(poi.get(field))
            if image_path:
                return image_path.read_bytes()
    return None



def add_brand_logo(slide, variant="mark", x=0.72, y=7.08, w=0.26, h=None):
    """Place the HZH logo if the packaged brand asset is available."""
    logo_path = LOGO_FULL_PATH if variant == "full" else LOGO_MARK_PATH
    if not logo_path.exists():
        return None
    try:
        image = Image.open(logo_path)
        image_w, image_h = image.size
        image.close()
        if h is None:
            h = w * image_h / image_w
        return slide.shapes.add_picture(str(logo_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    except Exception:
        return None

def add_background(slide):
    """Apply the proposal's light beige canvas and subtle top accent."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BEIGE

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = DEEP_GREEN
    accent.line.fill.background()


def add_title(slide, title, subtitle=None):
    """Add a consistent navy title and gold divider."""
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.48), Inches(11.9), Inches(0.55))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = clean_text(title)
    p.font.name = "Aptos Display"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    divider = slide.shapes.add_shape(1, Inches(0.72), Inches(1.18), Inches(2.4), Inches(0.04))
    divider.fill.solid()
    divider.fill.fore_color.rgb = GOLD
    divider.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.32), Inches(11.9), Inches(0.45))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        p.text = clean_text(subtitle)
        p.font.name = "Aptos"
        p.font.size = Pt(13)
        p.font.color.rgb = MUTED


def add_footer(slide, page_label="Private client proposal"):
    """Add a quiet footer to every slide."""
    line = slide.shapes.add_shape(1, Inches(0.72), Inches(7.02), Inches(11.9), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    add_brand_logo(slide, "mark", 0.72, 7.07, 0.24)
    box = slide.shapes.add_textbox(Inches(1.02), Inches(7.08), Inches(11.6), Inches(0.25))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = page_label
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED


def add_textbox(slide, text, x, y, w, h, size=14, color=TEXT, bold=False, min_size=10):
    """Add a plain editable text box with conservative wrapping and auto-shrink."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.clear()
    p = frame.paragraphs[0]
    p.text = clean_text(text)
    p.font.name = "Aptos"
    p.font.size = Pt(max(size, min_size))
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=15):
    """Add editable bullet points."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    safe_items = items or ["To be confirmed."]
    for index, item in enumerate(safe_items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = clean_text(item)
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
    return box


def add_centered_placeholder_text(shape, label):
    """Put centered placeholder copy inside an editable shape."""
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.18)
    frame.margin_right = Inches(0.18)
    frame.margin_top = Inches(0.12)
    frame.margin_bottom = Inches(0.12)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = clean_text(label)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = MUTED


def add_image_or_placeholder(slide, image_bytes_value, label, x, y, w, h):
    """Add a cropped image, or an editable dashed placeholder when empty."""
    left, top, width, height = Inches(x), Inches(y), Inches(w), Inches(h)

    if image_bytes_value:
        stream = BytesIO(image_bytes_value)
        picture = slide.shapes.add_picture(stream, left, top, width=width, height=height)

        img_w, img_h = image_size(image_bytes_value)
        image_ratio = img_w / img_h
        box_ratio = width / height

        # Crop to fill the requested image box while preserving the image ratio.
        if image_ratio > box_ratio:
            crop = (1 - (box_ratio / image_ratio)) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        else:
            crop = (1 - (image_ratio / box_ratio)) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop

        return picture

    placeholder = slide.shapes.add_shape(1, left, top, width, height)
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = PLACEHOLDER_FILL
    placeholder.line.color.rgb = PLACEHOLDER_LINE
    placeholder.line.width = Pt(1.6)
    placeholder.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if clean_text(label):
        add_centered_placeholder_text(placeholder, label)
    return placeholder


def adaptive_card_value_size(value, width, base_size=16, min_size=10):
    """Choose a readable card value size before PowerPoint auto-fit takes over."""
    text = clean_text(value)
    if not text:
        return base_size
    estimated_chars_per_line = max(int(width * 9), 8)
    lines_needed = max((len(text) + estimated_chars_per_line - 1) // estimated_chars_per_line, text.count("\n") + 1)
    if lines_needed >= 3 or len(text) > estimated_chars_per_line * 2:
        return max(min_size, base_size - 4)
    if lines_needed == 2 or len(text) > estimated_chars_per_line:
        return max(min_size, base_size - 2)
    return base_size


def format_route_for_cover(route_cities):
    """Format route text so multi-city routes can wrap cleanly on the cover."""
    cities = parse_route_cities(route_cities) or [clean_text(route_cities)]
    cities = [city for city in cities if city]
    if len(cities) >= 4:
        return "\n".join(cities)
    return " · ".join(cities) if cities else "To be confirmed"


def add_label_card(slide, label, value, x, y, w, h, value_size=16):
    """Add a small editable summary block with wrapped, self-contained text."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BEIGE
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1)

    add_textbox(slide, label.upper(), x + 0.18, y + 0.12, w - 0.36, 0.25, size=8.5, color=DEEP_GREEN, bold=True)
    safe_value_size = adaptive_card_value_size(value, w - 0.36, value_size, min_size=10)
    add_textbox(slide, value, x + 0.18, y + 0.43, w - 0.36, h - 0.48, size=safe_value_size, color=NAVY, bold=True, min_size=10)


def add_table(slide, headers, rows, x, y, w, h, font_size=11):
    """Create an editable PowerPoint table with premium proposal styling."""
    table_shape = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    table = table_shape.table

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BEIGE
        cell.text = clean_text(header)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = NAVY

    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else RGBColor(250, 246, 238)
            cell.text = clean_text(value)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = TEXT
            paragraph.alignment = PP_ALIGN.LEFT

    return table_shape


def add_cover(prs, title, group_size, days, nights, route_cities, intro, cover_image=None, day_pois=None):
    """Create the cover page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    add_brand_logo(slide, "full", 0.85, 0.58, 2.75)
    add_textbox(slide, "PRIVATE MEDICAL TOURISM PROPOSAL", 0.85, 1.25, 6.25, 0.35, size=12, color=DEEP_GREEN, bold=True)
    title_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.72), Inches(6.2), Inches(1.22))
    frame = title_box.text_frame
    frame.word_wrap = True
    frame.clear()
    p = frame.paragraphs[0]
    p.text = clean_text(title)
    p.font.name = "Aptos Display"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY

    divider = slide.shapes.add_shape(1, Inches(0.85), Inches(3.08), Inches(3.1), Inches(0.05))
    divider.fill.solid()
    divider.fill.fore_color.rgb = GOLD
    divider.line.fill.background()

    add_textbox(slide, intro, 0.85, 3.36, 6.05, 1.02, size=13)

    route_display = format_route_for_cover(route_cities)
    # Route is usually the longest cover metadata field, so it gets its own wide row.
    add_label_card(slide, "Group Size", f"{group_size} guests", 0.85, 4.75, 2.25, 0.88, value_size=15)
    add_label_card(slide, "Duration", f"{days} days / {nights} nights", 3.35, 4.75, 3.5, 0.88, value_size=14)
    add_label_card(slide, "Route", route_display, 0.85, 5.82, 6.0, 0.82, value_size=14)

    route_image = cover_image or image_bytes_from_poi_placeholder([poi for pois in (day_pois or []) for poi in pois])
    add_image_or_placeholder(slide, route_image, "", 7.75, 0.85, 4.75, 5.8)
    add_footer(slide)


def add_price_page(prs, pricing_result):
    """Create the package price page from the right-side pricing preview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Package Price", "Indicative pricing based on the selected itinerary and current quotation preview.")
    rows = [
        ["Client-facing Price", format_client_quote(pricing_result.get("client_quote"))],
    ]
    add_table(slide, ["Quotation Item", "Amount"], rows, 1.1, 2.1, 10.8, 0.9, font_size=15)
    add_textbox(slide, "Pricing is generated from the current itinerary, featured visits, service assumptions, and manual quotation adjustments. Final pricing is subject to confirmation.", 1.1, 4.1, 10.8, 0.7, size=13, color=MUTED)
    add_footer(slide)


def split_itinerary_for_overview_pages(itinerary_days, max_days_per_page=3):
    """Split itinerary overview into small pages so tables never overflow."""
    days = itinerary_records(itinerary_days)
    return [days[index:index + max_days_per_page] for index in range(0, len(days), max_days_per_page)]


def clean_overview_text(value):
    """Remove backend-only terms from customer-facing overview copy."""
    import re

    text = clean_text(value)
    blocked_terms = [
        "No POIs selected",
        "Selected POIs",
        "Image Placeholder",
        "Visual Highlight Image Placeholder",
        "placeholder",
        "undefined",
        "null",
    ]
    for term in blocked_terms:
        text = text.replace(term, "")
        text = text.replace(term.lower(), "")
    text = re.sub(r"\b[A-Za-z ]*POI[- ]?[A-Za-z0-9]+\b", "", text)
    text = re.sub(r"\s+", " ", text).strip(" -|,.;")
    return text


def short_overview_arrangement(row, max_chars=105):
    """Create a compact, client-facing arrangement for overview tables."""
    theme = clean_overview_text(row.get("Theme", ""))
    arrangement = clean_overview_text(row.get("Arrangement", ""))
    combined = f"{theme}. {arrangement}" if theme and theme.lower() not in arrangement.lower() else arrangement

    rules = [
        ("medical", "Medical checkup coordination with private transfer and recovery-friendly support."),
        ("consultation", "Medical consultation and recovery-friendly support with private transfer."),
        ("arrival", "Arrival reception, private transfer, hotel check-in, and welcome briefing."),
        ("departure", "Departure preparation with light arrangements and private transfer support."),
        ("business", "Curated business visit with local host and private transfer support."),
        ("innovation", "Curated innovation visit with local host and private transfer support."),
        ("wellness", "Wellness and light recovery experience at a comfortable pace."),
        ("culture", "Curated cultural visit and local lifestyle experience."),
    ]
    lower = combined.lower()
    for keyword, summary in rules:
        if keyword in lower:
            return summary

    if not arrangement:
        return "Curated local program with private transfer support."
    first_sentence = arrangement.split(".")[0].strip()
    if len(first_sentence) > max_chars:
        first_sentence = first_sentence[: max_chars - 1].rsplit(" ", 1)[0] + "."
    elif first_sentence and not first_sentence.endswith("."):
        first_sentence += "."
    return first_sentence or "Curated local program with private transfer support."


def add_itinerary_overview_slide(prs, page_days, page_index, total_pages):
    """Create one paginated overview slide with at most three days."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    day_numbers = [clean_text(row.get("Day")) for row in page_days if clean_text(row.get("Day"))]
    if day_numbers:
        subtitle = f"Day {day_numbers[0]}-{day_numbers[-1]} | Overview {page_index}/{total_pages}"
    else:
        subtitle = f"Overview {page_index}/{total_pages}"
    add_title(slide, "Itinerary Overview", subtitle)

    rows = []
    for row in page_days:
        rows.append(
            [
                clean_overview_text(row.get("Day", "")),
                clean_overview_text(row.get("City", "")),
                clean_overview_text(row.get("Theme", "")),
                short_overview_arrangement(row),
            ]
        )

    table_shape = add_table(slide, ["Day", "City", "Theme", "Key Arrangement"], rows, 0.75, 1.9, 11.85, 3.95, font_size=10)
    table = table_shape.table
    widths = [0.85, 1.75, 2.8, 6.45]
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for row in table.rows:
        row.height = Inches(0.92)
        for cell in row.cells:
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            cell.text_frame.margin_left = Inches(0.05)
            cell.text_frame.margin_right = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)

    add_footer(slide)


def add_itinerary_overview(prs, itinerary, day_pois=None):
    """Create paginated overview slides instead of forcing all days into one table."""
    pages = split_itinerary_for_overview_pages(itinerary, max_days_per_page=3)
    total_pages = max(len(pages), 1)
    for page_index, page_days in enumerate(pages or [[]], start=1):
        add_itinerary_overview_slide(prs, page_days, page_index, total_pages)


def add_daily_pages(prs, itinerary, day_images=None, day_pois=None):
    """Create one editable daily itinerary slide per row."""
    day_images = day_images or []
    day_pois = day_pois or []
    for index, row in enumerate(itinerary_records(itinerary)):
        day = clean_text(row.get("Day", ""))
        city = clean_text(row.get("City", ""))
        theme = clean_text(row.get("Theme", "Daily Program"))
        arrangement = client_safe_export_text(row.get("Arrangement", "To be confirmed.")) or generic_client_arrangement(index + 1, len(itinerary_records(itinerary)), city, theme)
        pois = day_pois[index] if index < len(day_pois) else []

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)
        add_title(slide, f"Day {day}: {city}", theme)
        add_label_card(slide, "City", city, 0.85, 1.85, 2.65, 1.0)
        add_label_card(slide, "Theme", theme, 3.8, 1.85, 2.95, 1.0)
        add_textbox(slide, "Arrangement", 0.85, 3.2, 2.5, 0.35, size=14, color=DEEP_GREEN, bold=True)
        add_textbox(slide, arrangement, 0.85, 3.62, 5.95, 0.95, size=13)
        highlights_text = poi_summary_text(pois)
        if highlights_text:
            add_textbox(slide, "Highlights", 0.85, 4.82, 2.5, 0.3, size=14, color=DEEP_GREEN, bold=True)
            add_textbox(slide, highlights_text, 0.85, 5.18, 5.95, 1.25, size=10.5)

        uploaded_image = day_images[index] if index < len(day_images) else None
        fallback_image = image_bytes_from_poi_placeholder(pois)
        add_image_or_placeholder(
            slide,
            uploaded_image or fallback_image,
            poi_image_label(day, pois),
            7.35,
            1.78,
            5.15,
            4.85,
        )
        add_footer(slide)


def add_highlights_page(prs, route_cities, visual_highlights=None):
    """Create a two-column visual highlights image grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Visual Highlights", "Client-facing touchpoints for the proposed experience.")

    highlights = visual_highlights or [
        {"Title": "Care Coordination", "Caption": "Appointment planning, local host support, and smooth communication."},
        {"Title": "Premium Stay", "Caption": "Comfortable hotels selected for recovery-friendly access and convenience."},
        {"Title": "Private Transfers", "Caption": "Airport reception and private ground transport throughout the route."},
        {"Title": "Route Experience", "Caption": f"Curated leisure moments across {clean_text(route_cities) or 'the selected cities'}."},
    ]
    positions = [(0.85, 1.75), (6.95, 1.75), (0.85, 4.22), (6.95, 4.22)]
    for index, ((x, y), highlight) in enumerate(zip(positions, highlights[:4]), start=1):
        title = clean_text(highlight.get("Title", f"Highlight {index}"))
        caption = clean_text(highlight.get("Caption", "Add a short client-facing caption."))
        image_bytes_value = highlight.get("Image")

        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.45), Inches(2.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BEIGE
        shape.line.color.rgb = GOLD
        shape.line.width = Pt(1)
        add_image_or_placeholder(slide, image_bytes_value, "", x + 0.2, y + 0.2, 2.1, 1.65)
        add_textbox(slide, title, x + 2.55, y + 0.25, 2.65, 0.42, size=16, color=NAVY, bold=True)
        add_textbox(slide, caption, x + 2.55, y + 0.82, 2.65, 0.8, size=11, color=TEXT)

    add_footer(slide)


def add_medical_addon_page(prs, implant_price):
    """Create the optional medical add-on page."""
    if not implant_price or implant_price <= 0:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Optional Medical Add-On", "Dental implant support can be added to the travel package.")
    add_label_card(slide, "Dental Implant Add-On", money(implant_price, "USD"), 0.9, 2.0, 4.0, 1.15)
    add_bullets(
        slide,
        [
            "Final medical suitability is subject to professional consultation.",
            "Treatment timing may vary based on diagnosis, imaging, and recovery needs.",
            "Clinical fees, inclusions, and exclusions should be confirmed before booking.",
        ],
        0.95,
        3.65,
        10.6,
        1.8,
        size=17,
    )
    add_footer(slide)


def add_includes_excludes_page(prs, includes, excludes):
    """Create the package includes and excludes page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Package Includes / Excludes", "Scope of services for the proposed package.")
    add_textbox(slide, "Package Includes", 0.9, 1.75, 5.2, 0.35, size=18, color=DEEP_GREEN, bold=True)
    add_textbox(slide, "Package Excludes", 7.0, 1.75, 5.2, 0.35, size=18, color=DEEP_GREEN, bold=True)
    add_bullets(slide, split_lines(includes), 0.9, 2.25, 5.3, 3.9, size=14)
    add_bullets(slide, split_lines(excludes), 7.0, 2.25, 5.3, 3.9, size=14)
    add_footer(slide)


def add_quotation_summary(prs, pricing_result):
    """Create the quotation summary page from the right-side pricing preview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title(slide, "Quotation Summary", "Summary for client review and final confirmation.")

    rows = [
        ["Package arrangement", "Based on the selected itinerary and service scope"],
        ["Hotel arrangement", format_optional_usd(pricing_result.get("hotel_cost"))],
        ["Transportation arrangement", format_optional_usd(pricing_result.get("transport_cost"))],
        ["Service arrangement", format_optional_usd(pricing_result.get("service_cost"))],
        ["Optional add-on", format_optional_usd(pricing_result.get("optional_addon_cost"))],
        ["Client-facing Price", format_client_quote(pricing_result.get("client_quote"))],
        ["Validity", "Subject to availability and final confirmation."],
    ]

    add_table(slide, ["Item", "Amount"], rows, 1.0, 1.6, 11.0, 4.4, font_size=11)
    add_textbox(slide, "Quotation values are based on the current itinerary, featured visits, service assumptions, and manual quotation inputs. Final pricing should be confirmed before client delivery.", 1.0, 6.25, 11.0, 0.45, size=11, color=MUTED)
    add_footer(slide)


def generate_pptx(
    title,
    group_size,
    days,
    nights,
    route_cities,
    pricing_result,
    intro,
    itinerary,
    includes,
    excludes,
    cover_image=None,
    day_images=None,
    visual_highlights=None,
    day_pois=None,
):
    """Build the editable proposal PowerPoint and return it as bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    pricing_result = pricing_result or calculate_pricing_preview(day_pois, None, None, None, None, None, "")

    add_cover(prs, title, group_size, days, nights, route_cities, intro, cover_image, day_pois)
    add_price_page(prs, pricing_result)
    add_itinerary_overview(prs, itinerary, day_pois)
    add_daily_pages(prs, itinerary, day_images, day_pois)
    add_highlights_page(prs, route_cities, visual_highlights)
    if pricing_result.get("optional_addon_cost"):
        add_medical_addon_page(prs, pricing_result.get("optional_addon_cost"))
    add_includes_excludes_page(prs, includes, excludes)
    add_quotation_summary(prs, pricing_result)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()



def parse_poi_images(poi):
    """Read the editable POI image list from JSON, falling back to legacy image fields."""
    images = []
    raw = clean_text(poi.get("images"))
    if raw:
        try:
            parsed = json.loads(raw)
            if isinstance(parsed, list):
                images = [item for item in parsed if isinstance(item, dict) and clean_text(item.get("url"))]
        except Exception:
            images = []

    if not images:
        legacy_url = clean_text(poi.get("image_path")) or clean_text(poi.get("image_url"))
        if legacy_url:
            images = [
                {
                    "id": f"IMG-{uuid4().hex[:8].upper()}",
                    "url": legacy_url,
                    "name": Path(legacy_url).name,
                    "isMain": True,
                    "uploadedAt": "",
                }
            ]

    if images and not any(bool(item.get("isMain")) for item in images):
        images[0]["isMain"] = True
    return images


def dump_poi_images(images):
    """Store image metadata as JSON in the CSV-backed POI database."""
    clean_images = []
    for image in images or []:
        url = clean_text(image.get("url"))
        if not url:
            continue
        clean_images.append(
            {
                "id": clean_text(image.get("id")) or f"IMG-{uuid4().hex[:8].upper()}",
                "url": url,
                "name": clean_text(image.get("name")) or Path(url).name,
                "isMain": bool(image.get("isMain")),
                "uploadedAt": clean_text(image.get("uploadedAt")) or time.strftime("%Y-%m-%d %H:%M:%S"),
            }
        )
    if clean_images and not any(image.get("isMain") for image in clean_images):
        clean_images[0]["isMain"] = True
    return json.dumps(clean_images, ensure_ascii=False)


def main_poi_image(images):
    """Return the image marked as main, or the first image."""
    if not images:
        return None
    for image in images:
        if image.get("isMain"):
            return image
    return images[0]


def image_summary_for_poi(poi):
    images = parse_poi_images(poi)
    if not images:
        return "暂无图片"
    return f"{len(images)}张" + ("（有主图）" if main_poi_image(images) else "")


def export_number(value):
    """Return numeric export values while leaving blanks empty."""
    number = optional_number(value)
    return "" if number is None else number


def build_poi_export_dataframe(poi_df):
    """Build a Chinese-header export table for the currently visible active POIs."""
    rows = []
    for _, row in poi_df.iterrows():
        poi = row.to_dict()
        images = parse_poi_images(poi)
        main_image = main_poi_image(images)
        image_urls = [clean_text(image.get("url")) for image in images if clean_text(image.get("url"))]
        rows.append(
            {
                "点位编号": clean_text(poi.get("poi_id")) or clean_text(poi.get("code")),
                "中文名称": clean_text(poi.get("name_cn")),
                "英文名称": clean_text(poi.get("name_en")),
                "点位类型": clean_text(poi.get("category")),
                "城市": clean_text(poi.get("city")),
                "区域": clean_text(poi.get("district")),
                "建议游览时长（小时）": export_number(poi.get("duration_hours")),
                "最低成本价": export_number(poi.get("price_min")),
                "最高成本价": export_number(poi.get("price_max")),
                "币种": clean_text(poi.get("currency")) or "CNY",
                "点位简介": clean_text(poi.get("description")) or clean_text(poi.get("description_cn")) or clean_text(poi.get("description_en")),
                "推荐理由": clean_text(poi.get("recommended_reason")),
                "游览亮点": clean_text(poi.get("highlights")),
                "适合人群": clean_text(poi.get("target_users")) or clean_text(poi.get("suitable_for")),
                "注意事项": clean_text(poi.get("notes")),
                "地址": clean_text(poi.get("address")),
                "开放时间": clean_text(poi.get("opening_hours")),
                "预约说明": clean_text(poi.get("reservation_info")),
                "内部备注": clean_text(poi.get("internal_remark")),
                "图片数量": len(images),
                "主图名称": clean_text(main_image.get("name")) if main_image else "",
                "主图地址": clean_text(main_image.get("url")) if main_image else "",
                "全部图片地址": "\n".join(image_urls),
                "标签": clean_text(poi.get("tags")),
                "状态": clean_text(poi.get("status")) or "active",
            }
        )
    return pd.DataFrame(rows)


def poi_export_xlsx_bytes(poi_df):
    """Return an Excel workbook for POI export."""
    export_df = build_poi_export_dataframe(poi_df)
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        export_df.to_excel(writer, index=False, sheet_name="POI点位列表")
        worksheet = writer.sheets["POI点位列表"]
        for column_cells in worksheet.columns:
            header = str(column_cells[0].value or "")
            max_length = max([len(str(cell.value or "")) for cell in column_cells] + [len(header)])
            worksheet.column_dimensions[column_cells[0].column_letter].width = min(max(max_length + 2, 12), 36)
        for row_cells in worksheet.iter_rows():
            for cell in row_cells:
                cell.alignment = cell.alignment.copy(wrap_text=True, vertical="top")
    buffer.seek(0)
    return buffer.getvalue()


def apply_images_to_poi_values(values, images):
    """Keep the new images JSON and old single-image fields in sync."""
    values["images"] = dump_poi_images(images)
    main_image = main_poi_image(images)
    if main_image:
        main_url = clean_text(main_image.get("url"))
        if main_url.lower().startswith(("http://", "https://")):
            values["image_url"] = main_url
            values["image_path"] = ""
        else:
            values["image_path"] = main_url
            values["image_url"] = ""
        values["image_notes"] = "主图"
    else:
        values["image_path"] = ""
        values["image_url"] = ""
        values["image_notes"] = ""
    return values



def table_cell(value, empty="暂无"):
    """Escape a value for display inside the POI HTML table."""
    text = clean_text(value)
    if not text:
        return f'<span class="muted">{html.escape(empty)}</span>'
    return html.escape(text)


def render_poi_list_table(poi_df, filtered_df):
    """Render a fixed-width, horizontally scrollable POI list with row actions."""
    if filtered_df.empty:
        st.info("当前筛选条件下暂无 POI。")
        return

    rows_html = []
    for _, row in filtered_df.iterrows():
        poi_id = clean_text(row.get("poi_id"))
        duration = optional_number(row.get("duration_hours"))
        duration_text = f"{duration:g}h" if duration is not None else "暂无"
        edit_href = f"?poi_action=edit&poi_id={quote(poi_id)}"
        delete_href = f"?poi_action=delete&poi_id={quote(poi_id)}"
        rows_html.append(
            "<tr>"
            f"<td class='hzh-poi-col-id nowrap'>{table_cell(poi_id)}</td>"
            f"<td class='hzh-poi-col-cn'>{table_cell(row.get('name_cn'))}</td>"
            f"<td class='hzh-poi-col-en'>{table_cell(row.get('name_en'))}</td>"
            f"<td class='hzh-poi-col-type'>{table_cell(row.get('category'))}</td>"
            f"<td class='hzh-poi-col-city nowrap'>{table_cell(row.get('city'))}</td>"
            f"<td class='hzh-poi-col-duration nowrap'>{html.escape(duration_text)}</td>"
            f"<td class='hzh-poi-col-price nowrap'>{html.escape(poi_price_range(row))}</td>"
            f"<td class='hzh-poi-col-image nowrap'>{html.escape(image_summary_for_poi(row))}</td>"
            "<td class='hzh-poi-col-action hzh-poi-action-cell'>"
            f"<a class='hzh-poi-action-btn edit' href='{edit_href}'>编辑</a>"
            f"<a class='hzh-poi-action-btn delete' href='{delete_href}'>删除</a>"
            "</td>"
            "</tr>"
        )

    table_html = """
    <div class="hzh-poi-list-note">列表页只负责查看、筛选、进入编辑和软删除。字段修改与图片管理请进入单个 POI 详情编辑页。</div>
    <div class="hzh-poi-table-wrapper">
      <table class="hzh-poi-table">
        <colgroup>
          <col class="hzh-poi-col-id" />
          <col class="hzh-poi-col-cn" />
          <col class="hzh-poi-col-en" />
          <col class="hzh-poi-col-type" />
          <col class="hzh-poi-col-city" />
          <col class="hzh-poi-col-duration" />
          <col class="hzh-poi-col-price" />
          <col class="hzh-poi-col-image" />
          <col class="hzh-poi-col-action" />
        </colgroup>
        <thead>
          <tr>
            <th>点位编号</th>
            <th>中文名称</th>
            <th>英文名称</th>
            <th>点位类型</th>
            <th>城市</th>
            <th>时长</th>
            <th>成本区间</th>
            <th>图片</th>
            <th>操作</th>
          </tr>
        </thead>
        <tbody>
          {rows}
        </tbody>
      </table>
    </div>
    """.format(rows="".join(rows_html))
    st.markdown(table_html, unsafe_allow_html=True)

    pending_delete_id = clean_text(st.session_state.get("pending_delete_poi_id"))
    if pending_delete_id:
        pending_rows = poi_df[poi_df["poi_id"].astype(str) == pending_delete_id]
        if not pending_rows.empty:
            pending_name = clean_text(pending_rows.iloc[0].get("name_cn")) or clean_text(pending_rows.iloc[0].get("name_en")) or pending_delete_id
            st.warning(f"确认删除 POI「{pending_name}」？删除后会从默认列表隐藏，但数据不会被物理删除。")
            c1, c2 = st.columns(2)
            if c1.button("确认删除", type="primary", key=f"confirm_delete_{pending_delete_id}"):
                save_pois(soft_delete_pois(poi_df, [pending_delete_id]))
                st.session_state.pending_delete_poi_id = ""
                st.query_params.clear()
                st.success("已软删除该 POI。")
                st.rerun()
            if c2.button("取消删除", key=f"cancel_delete_{pending_delete_id}"):
                st.session_state.pending_delete_poi_id = ""
                st.query_params.clear()
                st.rerun()


def render_poi_edit_page(poi_df, poi_id):
    """Render the second-level detail edit page for one POI."""
    selected_rows = poi_df[poi_df["poi_id"].astype(str) == str(poi_id)]
    if selected_rows.empty:
        st.error("未找到该 POI，可能已被删除或编号发生变化。")
        if st.button("返回列表"):
            st.session_state.main_nav = "POI 点位管理"
            st.session_state.poi_management_mode = "list"
            st.session_state.selected_poi_id = ""
            st.rerun()
        return

    poi = selected_rows.iloc[0].to_dict()
    poi_title = clean_text(poi.get("name_cn")) or clean_text(poi.get("name_en")) or clean_text(poi.get("poi_id"))

    title_cols = st.columns([0.38, 8], vertical_alignment="center")
    if title_cols[0].button("←", key=f"poi_back_icon_{poi_id}", help="返回 POI 列表"):
        st.session_state.poi_return_confirm = True
    title_cols[1].header(f"编辑 POI：{poi_title}")
    st.caption("这是单个 POI 的详情编辑页。保存、取消、图片管理都只在这里处理。")

    if st.session_state.get("poi_return_confirm"):
        st.warning("如果当前页面有未保存修改，返回列表后将放弃这些修改。")
        confirm_cols = st.columns([1.3, 1.3, 6])
        if confirm_cols[0].button("确认返回", type="primary", key=f"confirm_back_{poi_id}"):
            st.session_state.poi_return_confirm = False
            st.session_state.main_nav = "POI 点位管理"
            st.session_state.poi_management_mode = "list"
            st.session_state.selected_poi_id = ""
            st.rerun()
        if confirm_cols[1].button("继续编辑", key=f"keep_editing_{poi_id}"):
            st.session_state.poi_return_confirm = False
            st.rerun()

    images = parse_poi_images(poi)

    with st.form(f"poi_detail_form_{poi_id}"):
        st.subheader("基础信息")
        b1, b2, b3 = st.columns(3)
        with b1:
            code = st.text_input("点位编号", value=clean_text(poi.get("code")) or clean_text(poi.get("poi_id")))
            name_cn = st.text_input("中文名称", value=clean_text(poi.get("name_cn")))
            category = st.text_input("点位类型", value=clean_text(poi.get("category")))
            duration_hours = st.number_input("建议游览时长（小时）", min_value=0.0, value=float(optional_number(poi.get("duration_hours")) or 0), step=0.5)
        with b2:
            name_en = st.text_input("英文名称", value=clean_text(poi.get("name_en")))
            city = st.text_input("城市", value=clean_text(poi.get("city")))
            district = st.text_input("区域", value=clean_text(poi.get("district")))
            currency = st.text_input("币种", value=clean_text(poi.get("currency")) or "CNY")
        with b3:
            price_min = st.number_input("最低成本价", min_value=0.0, value=float(optional_number(poi.get("price_min")) or 0), step=100.0)
            price_max = st.number_input("最高成本价", min_value=0.0, value=float(optional_number(poi.get("price_max")) or 0), step=100.0)
            tags_value = st.text_input("标签", value=clean_text(poi.get("tags")))
            suitable_value = st.text_input("适合人群", value=clean_text(poi.get("suitable_for")))

        st.subheader("详情内容")
        d1, d2 = st.columns(2)
        with d1:
            description = st.text_area("点位简介 description", value=clean_text(poi.get("description")), height=90)
            recommended_reason = st.text_area("推荐理由 recommendedReason", value=clean_text(poi.get("recommended_reason")), height=90)
            highlights = st.text_area("游览亮点 highlights", value=clean_text(poi.get("highlights")), height=90)
            target_users = st.text_input("适合人群 targetUsers", value=clean_text(poi.get("target_users")) or clean_text(poi.get("suitable_for")))
            notes = st.text_area("注意事项 notes", value=clean_text(poi.get("notes")), height=80)
        with d2:
            address = st.text_input("地址 address", value=clean_text(poi.get("address")))
            opening_hours = st.text_input("开放时间 openingHours", value=clean_text(poi.get("opening_hours")))
            reservation_info = st.text_area("预约说明 reservationInfo", value=clean_text(poi.get("reservation_info")), height=80)
            internal_remark = st.text_area("内部备注 internalRemark", value=clean_text(poi.get("internal_remark")), height=80)
            description_cn = st.text_area("中文介绍", value=clean_text(poi.get("description_cn")), height=80)
            description_en = st.text_area("英文介绍", value=clean_text(poi.get("description_en")), height=80)

        st.subheader("图片管理")
        upload_images = st.file_uploader("上传图片", type=["jpg", "jpeg", "png", "webp"], accept_multiple_files=True, key=f"poi_detail_upload_{poi_id}")
        image_alt = st.text_input("图片英文说明", value=clean_text(poi.get("image_alt")))
        image_source = st.text_input("图片来源", value=clean_text(poi.get("image_source")))
        image_notes = st.text_input("图片备注", value=clean_text(poi.get("image_notes")))

        st.markdown("<div style='height: 0.75rem'></div>", unsafe_allow_html=True)
        action_cols = st.columns([1, 1, 6], gap="small")
        with action_cols[0]:
            save_clicked = st.form_submit_button("保存", type="primary")
        with action_cols[1]:
            cancel_clicked = st.form_submit_button("取消")

        if save_clicked:
            saved_images = list(images)
            for uploaded in upload_images or []:
                image_path = save_uploaded_poi_image(uploaded, poi_id, name_en or name_cn)
                saved_images.append(
                    {
                        "id": f"IMG-{uuid4().hex[:8].upper()}",
                        "url": image_path,
                        "name": clean_text(uploaded.name) or Path(image_path).name,
                        "isMain": not bool(saved_images),
                        "uploadedAt": time.strftime("%Y-%m-%d %H:%M:%S"),
                    }
                )
            values = {
                "id": clean_text(poi.get("id")) or poi_id,
                "poi_id": poi_id,
                "code": code or poi_id,
                "name_cn": name_cn,
                "name_en": name_en,
                "category": category,
                "city": city,
                "district": district,
                "duration_hours": duration_hours,
                "price_min": price_min,
                "price_max": price_max,
                "currency": currency or "CNY",
                "tags": tags_value,
                "suitable_for": suitable_value,
                "description": description,
                "description_cn": description_cn,
                "description_en": description_en,
                "recommended_reason": recommended_reason,
                "highlights": highlights,
                "target_users": target_users,
                "image_alt": image_alt,
                "image_source": image_source,
                "image_notes": image_notes,
                "address": address,
                "opening_hours": opening_hours,
                "reservation_info": reservation_info,
                "notes": notes,
                "internal_remark": internal_remark,
                "status": "active",
            }
            values = apply_images_to_poi_values(values, saved_images)
            save_pois(update_poi_row(poi_df, poi_id, values))
            st.session_state.poi_return_confirm = False
            st.session_state.main_nav = "POI 点位管理"
            st.session_state.poi_management_mode = "list"
            st.session_state.selected_poi_id = ""
            st.success("POI 已保存。")
            st.rerun()

        if cancel_clicked:
            st.session_state.poi_return_confirm = False
            st.session_state.main_nav = "POI 点位管理"
            st.session_state.poi_management_mode = "list"
            st.session_state.selected_poi_id = ""
            st.rerun()

    st.subheader("已有图片")
    if images:
        for image in images:
            image_id = clean_text(image.get("id"))
            image_url = clean_text(image.get("url"))
            c1, c2, c3, c4 = st.columns([1.2, 2.4, 1.0, 1.4])
            local_path = resolve_local_image_path(image_url)
            if local_path:
                c1.image(str(local_path), width=120)
            elif image_url.lower().startswith(("http://", "https://")):
                c1.image(image_url, width=120)
            else:
                c1.write("图片路径无效")
            c2.write(clean_text(image.get("name")) or Path(image_url).name)
            c3.write("主图" if image.get("isMain") else "普通图")
            if c4.button("设为主图", key=f"set_main_{poi_id}_{image_id}", disabled=bool(image.get("isMain"))):
                updated_images = []
                for item in images:
                    item = dict(item)
                    item["isMain"] = clean_text(item.get("id")) == image_id
                    updated_images.append(item)
                values = apply_images_to_poi_values({}, updated_images)
                save_pois(update_poi_row(poi_df, poi_id, values))
                st.rerun()
            if c4.button("删除图片", key=f"delete_image_{poi_id}_{image_id}"):
                updated_images = [item for item in images if clean_text(item.get("id")) != image_id]
                values = apply_images_to_poi_values({}, updated_images)
                save_pois(update_poi_row(poi_df, poi_id, values))
                st.rerun()
    else:
        st.info("该 POI 暂无图片。请在上方上传后点击保存。")


def render_poi_management(poi_df):
    """Render the two-level POI management module."""
    if "poi_management_mode" not in st.session_state:
        st.session_state.poi_management_mode = "list"
    if "selected_poi_id" not in st.session_state:
        st.session_state.selected_poi_id = ""
    if "pending_delete_poi_id" not in st.session_state:
        st.session_state.pending_delete_poi_id = ""

    poi_action = clean_text(st.query_params.get("poi_action", ""))
    query_poi_id = clean_text(st.query_params.get("poi_id", ""))
    if poi_action == "edit" and query_poi_id:
        st.session_state.poi_management_mode = "edit"
        st.session_state.selected_poi_id = query_poi_id
        st.session_state.pending_delete_poi_id = ""
        st.query_params.clear()
        st.rerun()
    if poi_action == "delete" and query_poi_id:
        st.session_state.pending_delete_poi_id = query_poi_id
        st.query_params.clear()
        st.rerun()

    if st.session_state.poi_management_mode == "edit" and st.session_state.selected_poi_id:
        render_poi_edit_page(poi_df, st.session_state.selected_poi_id)
        return

    st.header("POI 点位管理")
    st.caption("第一层：POI 列表页。每一行右侧提供编辑和删除入口；进入编辑后再管理详情和图片。")

    if POI_TEMPLATE_PATH.exists():
        st.download_button(
            "下载 POI CSV 模板",
            data=POI_TEMPLATE_PATH.read_bytes(),
            file_name="poi_template.csv",
            mime="text/csv",
        )

    with st.expander("导入 POI 数据", expanded=False):
        uploaded_poi_file = st.file_uploader(
            "上传 Excel / CSV",
            type=["xlsx", "xls", "csv"],
            help="支持 Excel 或 CSV。系统会自动识别字段、预览映射结果，并在确认后写入点位数据库。",
        )
        if uploaded_poi_file:
            try:
                import_result = analyze_poi_file(uploaded_poi_file.getvalue(), filename=uploaded_poi_file.name)
                mapping_rows = [
                    {
                        "系统字段": IMPORT_FIELD_LABELS.get(target, target),
                        "系统字段 key": target,
                        "原表字段": source,
                    }
                    for target, source in import_result.mapping.items()
                ]
                st.markdown("**字段映射结果**")
                st.dataframe(mapping_rows, use_container_width=True, hide_index=True)

                c1, c2, c3, c4 = st.columns(4)
                c1.metric("可导入点位", len(import_result.dataframe))
                c2.metric("原表重复行", import_result.duplicate_count)
                c3.metric("缺少名称", import_result.missing_required_count)
                c4.metric("价格为空", import_result.empty_price_count)

                st.markdown("**导入数据预览**")
                st.dataframe(display_poi_df(import_result.dataframe.head(20)), use_container_width=True, hide_index=True)

                import_mode = st.radio(
                    "导入方式",
                    ["只新增不覆盖", "覆盖同编号点位", "替换全部现有数据"],
                    horizontal=True,
                )
                if st.button("确认导入", type="primary"):
                    before_count = len(poi_df)
                    if import_mode == "替换全部现有数据":
                        final_df = import_result.dataframe
                    elif import_mode == "覆盖同编号点位":
                        final_df = merge_poi_data(poi_df, import_result.dataframe, overwrite=True)
                    else:
                        final_df = merge_poi_data(poi_df, import_result.dataframe, overwrite=False)
                    save_pois(final_df)
                    after_count = len(final_df)
                    added_count = max(after_count - before_count, 0)
                    st.success(
                        f"导入完成：成功处理 {len(import_result.dataframe)} 条；当前数据库共 {after_count} 条；新增 {added_count} 条；"
                        f"缺少名称 {import_result.missing_required_count} 条；价格为空 {import_result.empty_price_count} 条。"
                    )
                    st.rerun()
            except Exception as exc:
                st.error(f"导入失败：{exc}")

    st.subheader("POI 列表")
    management_df = active_pois_only(poi_df)
    filter_cols = st.columns(5)
    with filter_cols[0]:
        categories = st.multiselect("点位类型", sorted([value for value in management_df["category"].unique() if clean_text(value)]))
    with filter_cols[1]:
        cities = st.multiselect("城市", sorted([value for value in management_df["city"].unique() if clean_text(value)]))
    with filter_cols[2]:
        tags = st.multiselect("标签", unique_multi_values(management_df, "tags"))
    with filter_cols[3]:
        suitable_for = st.multiselect("适合人群", unique_multi_values(management_df, "suitable_for"))
    with filter_cols[4]:
        keyword = st.text_input("搜索名称")

    filtered = filter_pois(management_df, categories, cities, tags, suitable_for, keyword)
    if filtered.empty:
        st.info("当前没有可导出的 POI 数据。")
    else:
        export_file_name = f"POI点位列表_{time.strftime('%Y%m%d')}.xlsx"
        st.download_button(
            "导出 POI 列表",
            data=poi_export_xlsx_bytes(filtered),
            file_name=export_file_name,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            type="secondary",
        )
    render_poi_list_table(poi_df, filtered)

st.set_page_config(page_title="HZH 方案工具", layout="wide")
inject_app_theme()
ensure_poi_files()
poi_df = load_pois()
active_poi_df = active_pois_only(poi_df)

st.title("HZH 方案工具")
st.caption(f"界面语言：{interfaceLanguage}；导出文件语言：{outputLanguage}。用于生成英文客户方案 PPT，并管理每日行程 POI。")

if "main_nav" not in st.session_state:
    st.session_state.main_nav = "方案生成"

nav_query = clean_text(st.query_params.get("main_nav", ""))
if nav_query == "proposal":
    st.session_state.main_nav = "方案生成"
elif nav_query == "poi":
    st.session_state.main_nav = "POI 点位管理"

if clean_text(st.query_params.get("poi_action", "")) or st.session_state.get("poi_management_mode") == "edit":
    st.session_state.main_nav = "POI 点位管理"

main_nav = st.session_state.main_nav
proposal_active = "active" if main_nav == "方案生成" else ""
poi_active = "active" if main_nav == "POI 点位管理" else ""
st.markdown(
    f"""
    <nav class="hzh-app-tabs" aria-label="主功能导航">
      <a class="hzh-app-tab {proposal_active}" href="?main_nav=proposal">方案生成</a>
      <a class="hzh-app-tab {poi_active}" href="?main_nav=poi">POI 点位管理</a>
    </nav>
    """,
    unsafe_allow_html=True,
)

if main_nav == "方案生成":
    if "generationStatus" not in st.session_state:
        st.session_state.generationStatus = "idle"
    if "generated_intro" not in st.session_state:
        st.session_state.generated_intro = ""
    if "itinerary_rows" not in st.session_state:
        st.session_state.itinerary_rows = []
    if "generated_source_key" not in st.session_state:
        st.session_state.generated_source_key = ""
    if "generation_error" not in st.session_state:
        st.session_state.generation_error = ""
    if "draft_version" not in st.session_state:
        st.session_state.draft_version = 0
    if "day_poi_ids" not in st.session_state:
        st.session_state.day_poi_ids = []
    if "day_suggested_pois" not in st.session_state:
        st.session_state.day_suggested_pois = []
    if "route_pricing_preview" not in st.session_state:
        st.session_state.route_pricing_preview = {}
    if "pdf_draft" not in st.session_state:
        st.session_state.pdf_draft = None
    if "pdf_draft_ready" not in st.session_state:
        st.session_state.pdf_draft_ready = False
    if "pdf_overwrite_confirm" not in st.session_state:
        st.session_state.pdf_overwrite_confirm = False
    if "customer_pdf_bytes" not in st.session_state:
        st.session_state.customer_pdf_bytes = None
    if "pdf_draft_version" not in st.session_state:
        st.session_state.pdf_draft_version = 0
    if "current_proposal" not in st.session_state:
        st.session_state.current_proposal = None
    if "proposal_messages" not in st.session_state:
        st.session_state.proposal_messages = []
    if "proposal_versions" not in st.session_state:
        st.session_state.proposal_versions = []
    if "current_requirements" not in st.session_state:
        st.session_state.current_requirements = {}

    with st.sidebar:
        st.header("方案约束设置（可选）")
        st.caption("主输入入口在右侧客户需求对话区；这里仅用于查看识别结果和手动补充约束。")
        with st.expander("已识别客户需求摘要", expanded=True):
            st.markdown(requirement_summary(st.session_state.current_requirements))
        title = st.text_input("英文方案标题（可选）", "Premium Medical Tourism Proposal")
        group_size = st.number_input("出行人数（可选）", min_value=1, value=4, step=1)
        days = st.number_input("出行天数（可选）", min_value=1, value=7, step=1)
        nights = st.number_input("晚数（可选）", min_value=0, value=6, step=1)
        route_cities = st.text_input("城市范围（可选）", "Hangzhou, Yiwu, Shanghai")
        budget_min = st.text_input("客户预算下限（USD，可选）", "")
        budget_max = st.text_input("客户预算上限（USD，可选）", "")
        budget_notes = st.text_area("预算备注（可选）", "", height=80)
        client_quote_input = st.text_input(
            "客户展示报价（可选）",
            "",
            help="用于最终展示给客户的单一报价，例如 USD 559、To be confirmed 或 Based on final itinerary。",
        )
        customer_type = st.selectbox(
            "客户类型（可选）",
            ["综合客户", "康养客户", "商务客户", "家庭客户", "科技考察客户"],
        )
        preference_tags = st.multiselect(
            "行程偏好标签（可选）",
            ["医疗体检", "康养", "文化", "科技", "商务", "夜游", "美食", "轻体力", "室内", "高端", "亲子"],
        )
        intensity_level = st.selectbox("行程强度（可选）", ["标准", "轻松", "紧凑"])
        required_categories = st.multiselect(
            "必须包含的点位类型（可选）",
            ["医疗机构", "康养体验", "文化景点", "科技企业", "餐厅", "酒店", "夜游", "交通"],
        )
        excluded_categories = st.multiselect(
            "不希望包含的点位类型（可选）",
            ["夜游", "强步行", "户外", "商务参访", "科技企业", "餐厅"],
        )
        intro_input = ""
        generate_draft_clicked = False

    current_basics_key = proposal_basics_key(
        title,
        group_size,
        days,
        nights,
        route_cities,
        budget_min,
        budget_max,
        budget_notes,
        client_quote_input,
        intro_input,
        customer_type,
        preference_tags,
        intensity_level,
        required_categories,
        excluded_categories,
    )
    # 左侧现在只是可选约束设置，不再因为约束微调自动覆盖或标记当前方案过期。

    st.subheader("客户需求对话区")
    st.caption("这是方案生成的主入口。可以完全不填左侧，直接输入客户自然语言需求；生成后继续输入调整方向即可迭代当前方案。")
    with st.container(border=True):
        for message in st.session_state.proposal_messages[-6:]:
            role_label = "顾问" if message.get("role") == "user" else "系统"
            st.markdown(f"**{role_label}：** {message.get('content', '')}")
        dialogue_placeholder = "例如：客户是2个人，想做体检，也想参观杭州科技企业，7天，不要太累，预算中高。"
        dialogue_input = st.text_area("客户需求 / 调整指令", "", height=105, placeholder=dialogue_placeholder, key=f"proposal_dialogue_input_{st.session_state.draft_version}")
        col_chat_1, col_chat_2 = st.columns([1.2, 2.4])
        with col_chat_1:
            dialogue_action = "生成第一版方案" if not st.session_state.current_proposal else "根据当前方案继续调整"
            dialogue_clicked = st.button(dialogue_action, type="primary", disabled=not clean_text(dialogue_input))

    dialogue_filters = {
        "title": title,
        "customerCount": int(group_size),
        "days": int(days),
        "nights": int(nights),
        "routeCities": route_cities,
        "cities": parse_route_cities(route_cities),
        "budgetMin": budget_min,
        "budgetMax": budget_max,
        "budgetNote": budget_notes,
        "clientQuote": client_quote_input,
        "customerType": customer_type,
        "preferenceTags": preference_tags,
        "intensityLevel": intensity_level,
        "requiredCategories": required_categories,
        "excludedCategories": excluded_categories,
        "intro": intro_input,
    }
    if dialogue_clicked:
        try:
            st.session_state.proposal_messages.append({"role": "user", "content": dialogue_input})
            if st.session_state.current_proposal:
                sync_current_proposal_from_session()
                base_requirements = st.session_state.current_requirements or dialogue_filters
                extracted_requirements = infer_requirements_from_text(dialogue_input, base_requirements)
                st.session_state.current_requirements = {**base_requirements, **extracted_requirements}
                revised = revise_proposal(st.session_state.current_proposal, dialogue_input, st.session_state.current_requirements, active_poi_df)
                proposal_to_session_state(revised)
                append_proposal_version(f"迭代：{shorten_client_text(dialogue_input, 28)}")
                st.session_state.proposal_messages.append({"role": "assistant", "content": "已基于当前方案完成迭代修改，并同步到右侧方案草稿。"})
            else:
                extracted_requirements = infer_requirements_from_text(dialogue_input, dialogue_filters)
                st.session_state.current_requirements = extracted_requirements
                proposal, route_template = generate_initial_proposal(dialogue_input, extracted_requirements, active_poi_df)
                proposal_to_session_state(proposal, route_template.get("pricingPreview"))
                st.session_state.generated_source_key = current_basics_key
                append_proposal_version("v1 初始生成")
                st.session_state.proposal_messages.append({"role": "assistant", "content": "已生成第一版结构化方案草稿，可继续输入调整方向。"})
            st.rerun()
        except Exception as exc:
            st.session_state.generationStatus = "failed"
            st.session_state.generation_error = f"对话生成失败：{exc}"
            st.error(st.session_state.generation_error)

    if generate_draft_clicked:
        validation_errors = validate_proposal_basics(title, group_size, days, route_cities)
        if validation_errors:
            st.session_state.generationStatus = "failed"
            st.session_state.generation_error = "；".join(validation_errors)
        else:
            st.session_state.generationStatus = "generating"
            st.session_state.client_quote_input = client_quote_input
            st.session_state.pending_generation = {
                "intro": intro_input,
                "requirements": {
                    "title": title,
                    "customerCount": int(group_size),
                    "days": int(days),
                    "nights": int(nights),
                    "routeCities": route_cities,
                    "cities": parse_route_cities(route_cities),
                    "budgetMin": budget_min,
                    "budgetMax": budget_max,
                    "budgetNote": budget_notes,
                    "clientQuote": client_quote_input,
                    "customerType": customer_type,
                    "preferenceTags": preference_tags,
                    "intensityLevel": intensity_level,
                    "requiredCategories": required_categories,
                    "excludedCategories": excluded_categories,
                },
                "source_key": current_basics_key,
            }
            st.rerun()

    st.subheader("方案生成状态")
    status = st.session_state.generationStatus
    status_label = generationStatusLabels.get(status, "未知状态")
    if status == "completed":
        render_status_card("completed", f"方案生成状态：{status_label}", "方案初稿已生成，你可以继续输入调整方向，系统会基于当前方案进行修改。")
    elif status == "generating":
        render_status_card("generating", f"方案生成状态：{status_label}", "正在根据左侧信息生成英文方案初稿，请稍候……")
    elif status == "failed":
        render_status_card("failed", f"方案生成状态：{status_label}", st.session_state.generation_error or "生成失败，请检查左侧输入。")
    elif status == "stale":
        render_status_card("stale", "方案生成状态：方案约束已调整", "可继续使用当前方案，也可以在客户需求对话区输入调整方向进行迭代。")
    else:
        render_status_card("idle", "方案生成状态：未生成", "请输入客户需求，系统将生成第一版方案。")

    if status == "generating" and st.session_state.get("pending_generation"):
        time.sleep(0.7)
        try:
            pending = st.session_state.pending_generation
            route_template = generateRouteTemplateFromRequirements(pending["requirements"], active_poi_df)
            proposal = proposal_from_route_template(pending["requirements"], route_template, pending["intro"])
            proposal_to_session_state(proposal, route_template["pricingPreview"])
            append_proposal_version("表单生成方案初稿")
            st.session_state.generated_source_key = pending["source_key"]
            st.session_state.generationStatus = "completed"
            st.session_state.generation_error = ""
            st.session_state.pending_generation = None
            st.rerun()
        except Exception as exc:
            st.session_state.generationStatus = "failed"
            st.session_state.generation_error = f"生成失败：{exc}"
            st.session_state.pending_generation = None
            st.rerun()

    right_side_enabled = st.session_state.generationStatus in ["completed", "stale"] and bool(st.session_state.itinerary_rows)

    if right_side_enabled:
        generated_intro = st.text_area(
            "英文方案简介（可编辑预览）",
            value=st.session_state.generated_intro,
            height=110,
            key=f"generated_intro_editor_{st.session_state.draft_version}",
        )
        st.session_state.generated_intro = generated_intro

        st.subheader("英文行程安排")
        st.caption("可直接拖动表格最左侧手柄调整每天顺序；拖拽后系统会自动重排 Day 编号，并同步每日点位选择与导出顺序。")
        editor_rows = []
        for row_index, row in enumerate(st.session_state.itinerary_rows):
            editor_row = dict(row)
            editor_row["拖拽"] = "☰"
            editor_row["_row_id"] = row_index
            editor_rows.append(editor_row)
        try:
            from st_aggrid import AgGrid, GridOptionsBuilder, GridUpdateMode, DataReturnMode
            grid_df = pd.DataFrame(editor_rows)
            grid_builder = GridOptionsBuilder.from_dataframe(grid_df)
            grid_builder.configure_default_column(editable=True, resizable=True, wrapText=True, autoHeight=True)
            grid_builder.configure_column("拖拽", header_name="", rowDrag=True, editable=False, width=54, pinned="left")
            grid_builder.configure_column("Day", header_name="天数", editable=False, width=80)
            grid_builder.configure_column("City", header_name="英文城市", width=150)
            grid_builder.configure_column("Theme", header_name="英文主题", width=260)
            grid_builder.configure_column("Arrangement", header_name="英文行程安排", width=620)
            grid_builder.configure_column("_row_id", hide=True)
            grid_options = grid_builder.build()
            grid_options["rowDragManaged"] = True
            grid_options["rowDragEntireRow"] = True
            grid_options["animateRows"] = True
            grid_options["suppressMoveWhenRowDragging"] = False
            grid_response = AgGrid(
                grid_df,
                gridOptions=grid_options,
                update_mode=GridUpdateMode.MODEL_CHANGED,
                data_return_mode=DataReturnMode.FILTERED_AND_SORTED,
                fit_columns_on_grid_load=False,
                allow_unsafe_jscode=True,
                update_on=["rowDragEnd", "cellValueChanged"],
                height=min(620, 92 + 46 * max(len(editor_rows), 3)),
                key=f"itinerary_drag_grid_{st.session_state.draft_version}",
            )
            edited_records = itinerary_records(grid_response.get("data", editor_rows))
        except Exception as grid_exc:
            st.warning(f"拖拽表格组件未启用：{grid_exc}。请先安装依赖：pip install streamlit-aggrid")
            edited_records = itinerary_records(st.data_editor(
                editor_rows,
                num_rows="dynamic",
                use_container_width=True,
                column_order=["Day", "City", "Theme", "Arrangement"],
                column_config={
                    "Day": st.column_config.NumberColumn("天数", min_value=1, step=1, disabled=True),
                    "City": st.column_config.TextColumn("英文城市"),
                    "Theme": st.column_config.TextColumn("英文主题"),
                    "Arrangement": st.column_config.TextColumn("英文行程安排"),
                    "拖拽": None,
                    "_row_id": None,
                },
                key=f"itinerary_editor_fallback_{st.session_state.draft_version}",
            ))
        itinerary = []
        reordered_day_poi_ids = []
        reordered_suggested_pois = []
        for record in edited_records:
            original_index = int(optional_number(record.get("_row_id")) or 0)
            itinerary.append({
                "Day": record.get("Day"),
                "City": record.get("City"),
                "Theme": record.get("Theme"),
                "Arrangement": record.get("Arrangement"),
            })
            reordered_day_poi_ids.append(st.session_state.day_poi_ids[original_index] if original_index < len(st.session_state.day_poi_ids) else [])
            reordered_suggested_pois.append(st.session_state.day_suggested_pois[original_index] if original_index < len(st.session_state.day_suggested_pois) else [])
        old_order = list(range(len(st.session_state.itinerary_rows)))
        new_order = [int(optional_number(record.get("_row_id")) or 0) for record in edited_records]
        st.session_state.day_poi_ids = reordered_day_poi_ids
        st.session_state.day_suggested_pois = reordered_suggested_pois
        normalized_rows, normalized_day_poi_ids = normalize_itinerary_and_pois(itinerary_records(itinerary), st.session_state.day_poi_ids, active_poi_df)
        st.session_state.itinerary_rows = normalized_rows
        st.session_state.day_poi_ids = normalized_day_poi_ids
        while len(st.session_state.day_suggested_pois) < len(normalized_rows):
            idx = len(st.session_state.day_suggested_pois)
            row = normalized_rows[idx]
            st.session_state.day_suggested_pois.append(suggested_pois_for_day(row.get("City"), row.get("Theme")))
        st.session_state.day_suggested_pois = st.session_state.day_suggested_pois[:len(normalized_rows)]
        itinerary = normalized_rows
        if new_order != old_order and len(new_order) == len(old_order):
            sync_current_proposal_from_session()
            st.session_state.draft_version += 1
            st.toast("行程顺序已同步更新")
            st.rerun()
    else:
        generated_intro = ""
        itinerary = []
        st.subheader("英文行程安排")
        st.write("尚未生成方案初稿。")

    st.subheader("每日点位选择")
    poi_options = {poi_display_label(row): row["poi_id"] for _, row in active_poi_df.iterrows()}
    day_poi_ids = []
    if right_side_enabled:
        for index, row in enumerate(itinerary_records(itinerary), start=1):
            day_label = clean_text(row.get("Day", index)) or str(index)
            city_label = clean_text(row.get("City", ""))
            city_poi_options = {
                poi_display_label(poi_row): poi_row["poi_id"]
                for _, poi_row in active_poi_df.iterrows()
                if city_matches(poi_row.get("city"), city_label)
            }
            ordered_options = city_poi_options
            saved_ids = []
            if index - 1 < len(st.session_state.day_poi_ids):
                saved_ids = st.session_state.day_poi_ids[index - 1]
            default_labels = [
                label
                for label, poi_id in ordered_options.items()
                if str(poi_id) in {str(saved_id) for saved_id in saved_ids}
            ]
            labels = st.multiselect(
                f"第 {day_label} 天点位" + (f" - {city_label}" if city_label else ""),
                list(ordered_options.keys()),
                default=default_labels,
                key=f"day_pois_{st.session_state.draft_version}_{index}_{day_label}_{city_label}",
                help="系统会自动推荐初稿点位；这里仍可手动增加、删除或替换。",
            )
            selected_ids = [ordered_options[label] for label in labels]
            day_poi_ids.append(selected_ids)
            if index - 1 >= len(st.session_state.day_poi_ids):
                st.session_state.day_poi_ids.append(selected_ids)
            else:
                st.session_state.day_poi_ids[index - 1] = selected_ids
            if index - 1 >= len(st.session_state.day_suggested_pois):
                st.session_state.day_suggested_pois.append([])
            if selected_ids:
                st.session_state.day_suggested_pois[index - 1] = []
            elif not st.session_state.day_suggested_pois[index - 1]:
                st.session_state.day_suggested_pois[index - 1] = suggested_pois_for_day(city_label, row.get("Theme"))
            selected_db_pois = selected_poi_records(active_poi_df, selected_ids)
            badge_html = poi_source_badges_html(selected_db_pois, st.session_state.day_suggested_pois[index - 1])
            if badge_html:
                st.markdown("<div class='poi-source-legend'>蓝色=点位库已有；橙色=AI 临时建议，建议后续补入点位库。</div>" + badge_html, unsafe_allow_html=True)
            if not selected_ids:
                if ordered_options:
                    st.warning("当前天暂无自动入库点位，可从下拉框手动选择同城 POI；橙色标签为 AI 临时建议。")
                else:
                    st.warning("当前城市点位库暂无 POI；橙色标签为 AI 临时建议，请后续补充点位库。")
    else:
        st.write("方案初稿生成后，可在这里选择每日 POI。")

    st.session_state.itinerary_rows, day_poi_ids = normalize_itinerary_and_pois(itinerary_records(itinerary), day_poi_ids, active_poi_df)
    itinerary = st.session_state.itinerary_rows
    st.session_state.day_poi_ids = day_poi_ids
    while len(st.session_state.day_suggested_pois) < len(itinerary):
        idx = len(st.session_state.day_suggested_pois)
        row = itinerary[idx]
        st.session_state.day_suggested_pois.append(suggested_pois_for_day(row.get("City"), row.get("Theme")))
    st.session_state.day_suggested_pois = st.session_state.day_suggested_pois[:len(itinerary)]
    day_pois = [selected_poi_records(active_poi_df, ids) for ids in day_poi_ids]
    if right_side_enabled:
        sync_current_proposal_from_session()
    poi_min_cost, poi_max_cost = poi_cost_bounds(day_pois)
    st.caption(f"已选点位成本小计：{format_cost_range(poi_min_cost, poi_max_cost, 'CNY')}")

    with st.expander("预览已选点位详情", expanded=False):
        if day_pois:
            for index, pois in enumerate(day_pois, start=1):
                st.markdown(f"**第 {index} 天**")
                if pois:
                    preview_rows = [
                        {
                            "英文名称": clean_text(poi.get("name_en")),
                            "中文名称": clean_text(poi.get("name_cn")),
                            "建议时长": poi.get("duration_hours"),
                            "成本区间": poi_price_range(poi),
                            "图片占位符或路径": clean_text(poi.get("image_placeholder")),
                        }
                        for poi in pois
                    ]
                    st.dataframe(preview_rows, use_container_width=True, hide_index=True)
                else:
                    st.write("暂未选择点位。")
        else:
            st.write("暂未选择点位。")

    st.subheader("图片素材")
    cover_upload = st.file_uploader(
        "封面图片",
        type=["png", "jpg", "jpeg"],
        help="用于封面右侧的大图区域；不上传则保留可编辑图片占位符。",
    )

    st.markdown("**每日行程图片**")
    day_uploads = []
    if right_side_enabled:
        for index, row in enumerate(itinerary_records(itinerary), start=1):
            day_label = clean_text(row.get("Day", index)) or str(index)
            city_label = clean_text(row.get("City", ""))
            label = f"第 {day_label} 天图片"
            if city_label:
                label += f" - {city_label}"
            day_uploads.append(
                st.file_uploader(
                    label,
                    type=["png", "jpg", "jpeg"],
                    key=f"day_image_{index}_{day_label}_{city_label}",
                    help="用于该日 PPT 页面右侧大图区域；不上传则保留可编辑图片占位符。",
                )
            )
    else:
        st.write("方案初稿生成后，将按每日行程显示图片上传位。")

    st.markdown("**视觉亮点图片**")
    default_highlights = [
        {
            "Title": "Care Coordination",
            "Caption": "Appointment planning, local host support, and smooth communication.",
        },
        {
            "Title": "Premium Stay",
            "Caption": "Comfortable hotels selected for recovery-friendly access and convenience.",
        },
        {
            "Title": "Private Transfers",
            "Caption": "Airport reception and private ground transport throughout the route.",
        },
        {
            "Title": "Route Experience",
            "Caption": "Curated leisure moments across the selected cities.",
        },
    ]
    highlight_rows = st.data_editor(
        default_highlights,
        num_rows="dynamic",
        use_container_width=True,
        column_config={
            "Title": st.column_config.TextColumn("英文亮点标题"),
            "Caption": st.column_config.TextColumn("英文短说明"),
        },
        key="visual_highlight_rows",
    )
    highlight_uploads = st.file_uploader(
        "视觉亮点图片",
        type=["png", "jpg", "jpeg"],
        accept_multiple_files=True,
        help="图片会按顺序匹配上方亮点行；未上传的位置会保留可编辑图片占位符。",
    )
    highlight_uploads = highlight_uploads or []

    left, right = st.columns(2)
    with left:
        includes = st.text_area(
            "英文包含项目",
            "Private airport transfers\nSelected hotel accommodation\nDaily breakfast\nMedical appointment coordination\nLocal host and translation support\nCurated sightseeing arrangements",
            height=190,
        )
    with right:
        excludes = st.text_area(
            "英文不包含项目",
            "International flights\nVisa fees if applicable\nPersonal expenses\nTravel insurance\nMeals not listed in the itinerary\nMedical procedures not specified in the quotation",
            height=190,
        )

    st.subheader("价格核算 / 报价预览")
    if right_side_enabled:
        st.caption("左侧预算仅作为内部参考；下方报价根据右侧当前方案、已选 POI 和人工补充成本计算。")
        price_cols = st.columns(3)
        with price_cols[0]:
            hotel_cost = st.text_input("酒店成本小计（USD，可选）", "", key="hotel_cost_input")
            service_cost = st.text_input("服务成本小计（USD，可选）", "", key="service_cost_input")
        with price_cols[1]:
            transport_cost = st.text_input("交通成本小计（USD，可选）", "", key="transport_cost_input")
            optional_addon_cost = st.text_input("可选加项小计（USD，可选）", "", key="optional_addon_cost_input")
        with price_cols[2]:
            profit_multiplier = st.text_input("利润倍率 / 加价规则（可选）", "1.3", key="profit_multiplier_input")
            client_quote = st.text_input("客户展示报价（可选）", st.session_state.get("client_quote_input", ""), key="client_quote_input")

        pricing_result = calculate_pricing_preview(day_pois, hotel_cost, transport_cost, service_cost, optional_addon_cost, profit_multiplier, client_quote)
        pricing_rows = [
            {"项目": "基础成本小计", "金额": "待补充"},
            {"项目": "POI 成本小计", "金额": format_cost_range(pricing_result["poi_cost_min"], pricing_result["poi_cost_max"], "CNY")},
            {"项目": "酒店成本小计", "金额": format_optional_usd(pricing_result["hotel_cost"])},
            {"项目": "交通成本小计", "金额": format_optional_usd(pricing_result["transport_cost"])},
            {"项目": "服务成本小计", "金额": format_optional_usd(pricing_result["service_cost"])},
            {"项目": "可选加项小计", "金额": format_optional_usd(pricing_result["optional_addon_cost"])},
            {"项目": "总成本", "金额": format_cost_range(pricing_result["total_cost_min"] or 0, pricing_result["total_cost_max"] or 0, "CNY") if pricing_result["is_complete"] else "待补充"},
            {"项目": "客户展示报价", "金额": format_client_quote(pricing_result["client_quote"])},
        ]
        st.dataframe(pricing_rows, use_container_width=True, hide_index=True)
        missing_price_count = sum(
            1
            for pois in day_pois
            for poi in pois
            if optional_number(poi.get("price_min")) is None and optional_number(poi.get("price_max")) is None
        )
        if missing_price_count:
            st.warning(f"部分点位缺少成本价格，当前报价仅供内部参考。缺少价格点位：{missing_price_count} 个。")
        if not pricing_result["is_complete"]:
            st.info("当前方案成本数据不完整，请补充 POI / 酒店 / 医疗项目价格后再生成完整报价。")
    else:
        pricing_result = calculate_pricing_preview([], None, None, None, None, None, "")
        st.write("方案初稿生成后，可在这里进行价格核算和报价预览。")


    st.subheader("客户宣传页 PDF")
    st.caption("PDF 使用两阶段流程：先生成可编辑内容草稿，确认修改后再生成最终客户宣传页 PDF。")
    pdf_ready = st.session_state.generationStatus == "completed" and bool(itinerary_records(itinerary))
    pdf_cols = st.columns([1, 1, 2])
    with pdf_cols[0]:
        create_pdf_draft_clicked = st.button("生成 PDF 内容草稿", type="secondary", disabled=not pdf_ready)
    with pdf_cols[1]:
        regenerate_pdf_draft_clicked = False
        if st.session_state.get("pdf_draft_ready"):
            regenerate_pdf_draft_clicked = st.button("重新生成草稿", type="secondary")
    if not pdf_ready:
        st.info("请先生成方案后再生成 PDF 内容草稿。")

    def create_customer_pdf_draft(keep_client_name=False):
        existing_client_name = ""
        if keep_client_name and st.session_state.get("pdf_draft"):
            existing_client_name = st.session_state.pdf_draft.get("cover", {}).get("clientName", "")
        return build_pdf_draft_from_proposal(
            title=title,
            client_name=existing_client_name,
            group_size=group_size,
            days=days,
            nights=nights,
            route_cities=route_cities,
            intro=generated_intro,
            itinerary=itinerary_records(itinerary),
            day_pois=day_pois,
            includes=includes,
            excludes=excludes,
            pricing_result=pricing_result,
            customer_type=customer_type,
        )

    if create_pdf_draft_clicked:
        if st.session_state.get("pdf_draft_ready") and st.session_state.get("pdf_draft"):
            st.session_state.pdf_overwrite_confirm = True
        else:
            st.session_state.pdf_draft = create_customer_pdf_draft()
            st.session_state.pdf_draft_ready = True
            st.session_state.customer_pdf_bytes = None
            st.session_state.pdf_draft_version += 1
            st.success("PDF 内容草稿已生成，可在下方编辑后再生成最终 PDF。")

    if regenerate_pdf_draft_clicked:
        st.session_state.pdf_overwrite_confirm = True

    if st.session_state.get("pdf_overwrite_confirm"):
        st.warning("当前已有 PDF 草稿内容，重新生成会覆盖已编辑内容，是否继续？")
        confirm_cols = st.columns([1, 1, 3])
        with confirm_cols[0]:
            if st.button("确认覆盖草稿", type="primary"):
                st.session_state.pdf_draft = create_customer_pdf_draft(keep_client_name=True)
                st.session_state.pdf_draft_ready = True
                st.session_state.pdf_overwrite_confirm = False
                st.session_state.customer_pdf_bytes = None
                st.session_state.pdf_draft_version += 1
                st.success("PDF 内容草稿已重新生成。")
                st.rerun()
        with confirm_cols[1]:
            if st.button("保留当前草稿"):
                st.session_state.pdf_overwrite_confirm = False
                st.rerun()

    if st.session_state.get("pdf_draft_ready") and st.session_state.get("pdf_draft"):
        draft = st.session_state.pdf_draft
        draft_version = st.session_state.pdf_draft_version
        st.markdown("### 客户宣传页内容编辑区")
        edit_tabs = st.tabs([
            "封面信息",
            "方案概览",
            "顾问推荐意见",
            "每日行程内容",
            "Included / Not Included",
            "Pricing Summary",
            "Next Steps",
        ])

        with edit_tabs[0]:
            cover = draft.setdefault("cover", {})
            cover["title"] = st.text_input("PDF 标题", cover.get("title", ""), key=f"pdf_cover_title_{draft_version}")
            cover["clientName"] = st.text_input("客户姓名", cover.get("clientName", "Valued Client"), key=f"pdf_cover_client_{draft_version}")
            cover["publishedDate"] = st.text_input("生成日期", cover.get("publishedDate", ""), key=f"pdf_cover_date_{draft_version}")
            cover["validity"] = st.text_input("有效期", cover.get("validity", "30 days"), key=f"pdf_cover_validity_{draft_version}")
            cover["heroTitle"] = st.text_input("主标题", cover.get("heroTitle", ""), key=f"pdf_cover_hero_{draft_version}")
            cover["subtitle"] = st.text_area("副标题 / 宣传页说明", cover.get("subtitle", cover.get("introText", "")), height=90, key=f"pdf_cover_subtitle_{draft_version}")
            cover["introText"] = st.text_area("简短介绍文案", cover.get("introText", ""), height=130, key=f"pdf_cover_intro_{draft_version}")
            cover["quotedPrice"] = st.text_input("客户展示报价", cover.get("quotedPrice", "To be confirmed"), key=f"pdf_cover_quote_{draft_version}")
            cover["heroImagePath"] = st.text_input("首页主视觉图片路径 / 自动选图", cover.get("heroImagePath", ""), key=f"pdf_cover_image_{draft_version}")

        with edit_tabs[1]:
            overview = draft.setdefault("overview", {})
            overview["tripType"] = st.text_input("Trip Type", overview.get("tripType", ""), key=f"pdf_overview_trip_{draft_version}")
            overview["duration"] = st.text_input("Duration", overview.get("duration", ""), key=f"pdf_overview_duration_{draft_version}")
            overview["medicalCore"] = st.text_input("Medical Core", overview.get("medicalCore", ""), key=f"pdf_overview_medical_{draft_version}")
            overview["service"] = st.text_input("Service", overview.get("service", ""), key=f"pdf_overview_service_{draft_version}")
            overview["summary"] = st.text_area("方案摘要", overview.get("summary", ""), height=140, key=f"pdf_overview_summary_{draft_version}")

        with edit_tabs[2]:
            recommendation = draft.setdefault("recommendation", {})
            recommendation["title"] = st.text_input("Consultant Recommendation 标题", recommendation.get("title", ""), key=f"pdf_rec_title_{draft_version}")
            recommendation["body"] = st.text_area("推荐说明正文", recommendation.get("body", ""), height=140, key=f"pdf_rec_body_{draft_version}")
            recommendation["medicalCore"] = st.text_input("医疗核心说明", recommendation.get("medicalCore", ""), key=f"pdf_rec_medical_{draft_version}")
            recommendation["highlights"] = st.text_area("体检套餐亮点", recommendation.get("highlights", ""), height=110, key=f"pdf_rec_highlights_{draft_version}")
            recommendation["targetUsers"] = st.text_area("适合人群", recommendation.get("targetUsers", ""), height=90, key=f"pdf_rec_target_{draft_version}")
            recommendation["serviceNote"] = st.text_area("服务说明", recommendation.get("serviceNote", ""), height=110, key=f"pdf_rec_service_{draft_version}")
            recommendation["imagePath"] = st.text_input("医疗/顾问推荐图片路径 / 自动选图", recommendation.get("imagePath", ""), key=f"pdf_rec_image_{draft_version}")

        with edit_tabs[3]:
            itinerary_draft = draft.setdefault("itinerary", [])
            if not itinerary_draft:
                st.info("当前没有可编辑的每日行程。")
            for day_index, day in enumerate(itinerary_draft):
                with st.expander(f"第 {day.get('day', day_index + 1)} 天 - {day.get('title', '')}", expanded=day_index == 0):
                    day["dayTitle"] = st.text_input("Day 标题", day.get("dayTitle", f"Day {day_index + 1}"), key=f"pdf_day_title_{draft_version}_{day_index}")
                    day["title"] = st.text_input("当日英文标题", day.get("title", ""), key=f"pdf_day_headline_{draft_version}_{day_index}")
                    day["description"] = st.text_area("当日行程描述", day.get("description", ""), height=120, key=f"pdf_day_desc_{draft_version}_{day_index}")
                    day["pois"] = st.text_area("当日包含的 POI（每行一条）", day.get("pois", ""), height=90, key=f"pdf_day_pois_{draft_version}_{day_index}")
                    day["tags"] = st.text_input("当日标签", day.get("tags", ""), key=f"pdf_day_tags_{draft_version}_{day_index}")
                    day["imagePath"] = st.text_input("当日图片路径 / 自动选图", day.get("imagePath", ""), key=f"pdf_day_image_{draft_version}_{day_index}")
                    uploaded_pdf_day_image = st.file_uploader(
                        "上传当天图片（可选）",
                        type=["png", "jpg", "jpeg", "webp"],
                        key=f"pdf_day_upload_{draft_version}_{day_index}",
                    )
                    if uploaded_pdf_day_image:
                        POI_IMAGE_DIR.mkdir(parents=True, exist_ok=True)
                        suffix = Path(uploaded_pdf_day_image.name).suffix.lower() or ".png"
                        image_target = POI_IMAGE_DIR / f"pdf_day_{draft_version}_{day_index}_{uuid4().hex[:8]}{suffix}"
                        image_target.write_bytes(uploaded_pdf_day_image.getbuffer())
                        day["imagePath"] = str(image_target)
                    if day.get("imagePath"):
                        st.caption(f"当前图片：{day.get('imagePath')}")

        with edit_tabs[4]:
            draft["included"] = st.text_area("Included 内容（每行一条）", draft.get("included", ""), height=180, key=f"pdf_included_{draft_version}")
            draft["notIncluded"] = st.text_area("Not Included / To Confirm 内容（每行一条）", draft.get("notIncluded", ""), height=180, key=f"pdf_not_included_{draft_version}")

        with edit_tabs[5]:
            pricing_draft = draft.setdefault("pricing", {})
            pricing_draft["clientFacingPrice"] = st.text_input("Client-facing Price", pricing_draft.get("clientFacingPrice", "To be confirmed"), key=f"pdf_price_quote_{draft_version}")
            pricing_draft["pricingNote"] = st.text_area("价格说明", pricing_draft.get("pricingNote", ""), height=120, key=f"pdf_price_note_{draft_version}")
            pricing_draft["finalNote"] = st.text_area("最终价格确认说明", pricing_draft.get("finalNote", ""), height=100, key=f"pdf_price_final_{draft_version}")

        with edit_tabs[6]:
            next_steps = draft.setdefault("nextSteps", {})
            next_steps["items"] = st.text_area("客户需要确认的信息（每行一条）", next_steps.get("items", ""), height=170, key=f"pdf_next_items_{draft_version}")
            next_steps["cta"] = st.text_input("下一步说明", next_steps.get("cta", "Contact HZH consultant to confirm availability"), key=f"pdf_next_cta_{draft_version}")
            next_steps["note"] = st.text_area("联系或确认说明", next_steps.get("note", ""), height=110, key=f"pdf_next_note_{draft_version}")

        st.session_state.pdf_draft = draft
        pdf_action_cols = st.columns([1, 1, 2])
        with pdf_action_cols[0]:
            if st.button("生成最终客户宣传页 PDF", type="primary"):
                try:
                    st.session_state.customer_pdf_bytes = generate_customer_pdf(
                        st.session_state.pdf_draft,
                        logo_path=str(LOGO_FULL_PATH) if LOGO_FULL_PATH.exists() else None,
                    )
                    st.success("最终客户宣传页 PDF 已生成。")
                except ImportError:
                    st.error("当前环境缺少 PDF 依赖 ReportLab。请先运行 pip install -r requirements.txt 后再生成 PDF。")
                except Exception as exc:
                    st.error(f"生成客户宣传页 PDF 失败：{exc}")
        if st.session_state.get("customer_pdf_bytes"):
            st.download_button(
                "下载客户宣传页 PDF",
                data=st.session_state.customer_pdf_bytes,
                file_name=customer_pdf_filename(st.session_state.pdf_draft),
                mime="application/pdf",
            )

    export_ready = st.session_state.generationStatus == "completed" and bool(itinerary_records(itinerary))
    if not export_ready:
        st.warning("请先在客户需求对话区生成第一版方案，完成后再导出 PPT。")

    if st.button("生成英文方案 PPT", type="primary", disabled=not export_ready):
        visual_highlights = []
        for index, row in enumerate(itinerary_records(highlight_rows)):
            visual_highlights.append(
                {
                    "Title": row.get("Title", f"Highlight {index + 1}"),
                    "Caption": row.get("Caption", "Add a short client-facing caption."),
                    "Image": uploaded_file_bytes(highlight_uploads[index]) if index < len(highlight_uploads) else None,
                }
            )

        pptx = generate_pptx(
            title=title,
            group_size=group_size,
            days=days,
            nights=nights,
            route_cities=route_cities,
            pricing_result=pricing_result,
            intro=generated_intro,
            itinerary=itinerary,
            includes=includes,
            excludes=excludes,
            cover_image=uploaded_file_bytes(cover_upload),
            day_images=[uploaded_file_bytes(day_upload) for day_upload in day_uploads],
            visual_highlights=visual_highlights,
            day_pois=day_pois,
        )
        safe_name = title.lower().replace(" ", "_").replace("/", "_")
        st.success("英文方案 PPT 已生成。")
        st.download_button(
            "下载英文方案 PPT",
            data=pptx,
            file_name=f"{safe_name}_proposal.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

if main_nav == "POI 点位管理":
    render_poi_management(poi_df)
